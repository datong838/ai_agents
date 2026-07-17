# -*- coding: utf-8 -*-
"""Generate Palantir Foundry / AIP / Ontology deep-dive PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = r"c:\work\projects\wchat\docs\ref\Palantir-Foundry-AIP-Ontology-深度解析.pptx"

# Colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)


def add_bullets(slide, items, left=0.5, top=1.35, width=12.3, height=5.8, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(font_size if p.level == 0 else font_size - 2)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def add_mono_block(slide, text, left=0.5, top=1.35, width=12.3, height=5.8, size=11):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "Palantir Foundry & AIP"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.0))
    sp = s.text_frame.paragraphs[0]
    sp.text = "Ontology 可执行语义内核 · 企业 AI 编排治理层\n深度技术解析"
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    f = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11), Inches(0.5))
    fp = f.text_frame.paragraphs[0]
    fp.text = "谛听技术调研 · 2026-07-12 · 基于 Palantir 官方文档与公开资料整理"
    fp.font.size = Pt(12)
    fp.font.color.rgb = GRAY


def slide_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_title_bar(slide, "目录")
    add_bullets(slide, [
        "一、Palantir 产品矩阵与 Foundry 定位",
        "二、Ontology：可执行语义内核（四层栈 · 三大区 · 微服务）",
        "三、OSv1 Phonograph → OSv2 架构演进与 Search Around",
        "四、AIP：k-LLM · Logic · Agent Studio · Assist · Evals",
        "五、Apollo：Hub-Spoke 约束求解与自主部署操作系统",
        "六、Foundry 300+ 微服务编排推演 · AIP×Ontology 时序",
        "七、技术护城河总结",
    ], font_size=18)


def new_content_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT)
    add_title_bar(slide, title, subtitle)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_toc(prs)

    # --- Part 1 Foundry ---
    s = new_content_slide(prs, "Palantir 产品矩阵", "Gotham · Foundry · AIP · Apollo")
    add_bullets(s, [
        "Gotham：政府/国防/情报场景的数据融合与调查分析（Palantir 起源产品）",
        "Foundry：面向商业企业的中央数据运营平台，被称为「企业数据操作系统」",
        "AIP（2023）：生成式 AI 操作系统，为 Gotham/Foundry 接入 LLM 与 Agent 的安全编排层",
        "Apollo：跨环境软件交付与运维（含断网/VPC 战术部署），支撑 Foundry/AIP 私有化",
        "核心范式：数据湖（Dataset）→ Ontology（语义+动作）→ 应用/AI（Workshop/AIP）",
    ])

    s = new_content_slide(prs, "Foundry：企业数据操作系统", "定位与核心能力")
    add_bullets(s, [
        "定位：解决企业数据孤岛，统一集成、治理、分析、运营决策",
        "200+ 预建连接器：ERP / CRM / IoT / 数据库 / 文件等数据源接入",
        "数据管道：Pipeline Builder、Transforms（Python/SQL/Spark）清洗产出 Dataset",
        "存储底座：Multimodal Data Plane（MMDP），基于 Iceberg/Parquet 等开放格式",
        "分析工具：Contour（表格）、Quiver（时序/图表）、Workshop（低代码运营应用）",
        "治理能力：数据血缘、数字孪生、权限标记（Markings）、全流程审计",
        "典型场景：供应链、制造、医疗、能源、金融运营决策",
    ])

    # --- Part 2 Ontology ---
    s = new_content_slide(prs, "Ontology 重新定义", "不是知识图谱，而是 Executable Semantic Kernel")
    add_bullets(s, [
        "传统误解：Ontology ≠ 单纯知识图谱 ≠ 元数据目录",
        "正确定位：Foundry/AIP 的「可执行语义内核」（Executable Semantic Kernel）",
        "物理位置：介于 MMDP 数据湖 与 Applications/AIP 之间",
        "解决三件事：",
        ("业务语义建模（Semantic）— 世界是什么", 1),
        ("受控操作执行（Kinetic）— 能做什么", 1),
        ("动态治理（Governance）— 谁能在什么条件下做", 1),
        "护城河：把「数据表」升维为「带权限、带逻辑、带动作的业务对象世界」",
    ])

    s = new_content_slide(prs, "Ontology 四层栈架构")
    add_mono_block(s, """┌─────────────────────────────────────────────────────────────┐
│  L4  应用与 AI 层（Consumer）                                 │
│      Workshop · Slate · Quiver · Vertex · OSDK · AIP Agent    │
├─────────────────────────────────────────────────────────────┤
│  L3  本体核心引擎（Ontology Core）                             │
│      OMS · Object DB(OSv2) · OSS · Actions · Functions        │
├─────────────────────────────────────────────────────────────┤
│  L2  本体映射层（Mapping）— Object Data Funnel                │
│      Backing Dataset 列 → Object 属性/主键/Link · CDC 增量索引  │
├─────────────────────────────────────────────────────────────┤
│  L1  物理数据层（Raw Data）— Iceberg/Parquet Dataset          │
│      Connectors → Pipeline/Transforms → Backing Dataset       │
└─────────────────────────────────────────────────────────────┘
原则：底层数据不动；Funnel 构建物化对象索引；应用不直接碰底层表""")

    s = new_content_slide(prs, "三大并列区", "Semantic · Kinetic · Governance（运行时原子裁决）")
    add_bullets(s, [
        "Semantic 区（名词）：Object Type · Property · Link Type · Interface",
        "  → 定义消歧实体与关系（关系可带属性与时间维度）",
        "Kinetic 区（动词）：Action Type · Functions on Objects",
        "  → Action = 结构化状态迁移；Function = 无副作用计算或可调用逻辑",
        "Governance 区（治理）：Security Markings · Restrictions · Rules",
        "  → 权限内生于类型系统；按用户/Agent/上下文实时计算行列级策略",
        "设计精髓：三区并列而非堆叠，读写时在 Ontology 引擎内一次性裁决",
    ])

    s = new_content_slide(prs, "核心后端微服务", "解耦微服务协同")
    add_bullets(s, [
        "OMS（Ontology Metadata Service）：Schema 注册表，Object/Link/Action/Interface 版本化",
        "Object DB：OSv1 Phonograph（遗留）/ OSv2（当前主流，数十亿对象、毫秒点查）",
        "Object Data Funnel：批/流/CDC + 用户 Edit → 索引到 Object DB，最终一致",
        "OSS（Object Set Service）：统一读网关—过滤、搜索、聚合、Link Traversal",
        "Actions Service：唯一受控写入口—RBAC+行列策略+业务规则+审计+可选回写 ERP",
        "Functions on Objects：TS/Python 服务端逻辑，供 Actions/Workshop/AIP Agent 作 Tool",
    ], font_size=15)

    s = new_content_slide(prs, "OSv1 Phonograph → OSv2 演进", "官方：OSv1 计划 2026-06-30 后不可用")
    add_bullets(s, [
        "OSv1（Phonograph）— 遗留架构：",
        ("索引与查询子系统耦合；暴露大量底层 DB API", 1),
        ("Datasource 注册 → 全量/增量 reindex → Phonograph 查询", 1),
        ("用户 Edit 依赖 writeback dataset", 1),
        "OSv2 — 下一代canonical store（从第一性原理重构）：",
        ("索引子系统与查询子系统解耦 → 水平扩展", 1),
        ("全对象类型支持增量索引（CDC）", 1),
        ("Actions 经 Funnel 统一写入；materialized dataset 可选（非强制 writeback）", 1),
        ("Search Around / 大规模聚合 → Spark 并行计算层", 1),
        ("更严格数据校验、列级权限、确定性行为", 1),
        "迁移：Ontology Manager 支持逐类型迁移，soak period 最长 14 天可回滚",
    ], font_size=14)

    s = new_content_slide(prs, "Search Around 分布式执行", "OSv2 + Vertex Graph Functions")
    add_bullets(s, [
        "Search Around：从 Object Set A 沿 Link 遍历到 Object Set B 的图扩展查询",
        "OSS 下发请求 → OSv2 索引裁剪 → 权限过滤 → 返回 Object Set",
        "规模分层执行（官方 Query Compute 文档）：",
        ("常规：<100K 对象，索引后端低延迟路径", 1),
        (">100K 对象：按需启动 Spark 容器并行执行 Search Around", 1),
        (">10K 对象单次 writeback：同样可走 Spark 并行", 1),
        "Vertex Generate Graph Functions：",
        ("TS Function 返回 IGraphSearchAroundResultV1", 1),
        ("定义 directEdges / intermediateEdges / 布局，供 Vertex 可视化", 1),
        ("Workshop → Vertex 模板 → Search Around Function 参数化调用", 1),
        "工程建议：先过滤缩小 Object Set，再执行 Search Around / 聚合",
    ], font_size=14)

    s = new_content_slide(prs, "关键技术特征", "MDO · OSDK · AI Grounding · Interfaces")
    add_bullets(s, [
        "虚拟视图 + 物化索引：Backing Dataset 在 Iceberg 不动，Funnel 物化到 OSv2",
        "MDO（Multi-Datasource Objects）：一 Object Type 横向合并多源（列/行级 MDO）",
        "强类型 OSDK：Conjure RPC/IDL → TS/Java/Python SDK，应用像操作本地对象",
        "Interfaces 多态：IApprovable 等共享能力，上层统一交互无需关心具体类型",
        "AI Grounding 边界：AIP Agent 不能直接访问 Dataset，只能经 Ontology Tool：",
        ("Query Objects · Call Function · Apply Action（权限裁剪后视图）", 1),
        "Decision Lineage：谁/何时/基于哪版数据/经哪个 App 做了什么变更—完整审计",
    ], font_size=14)

    s = new_content_slide(prs, "Ontology 数据流闭环")
    add_mono_block(s, """Source Systems → Connectors → Foundry Datasets → Transforms
        ↓
   Backing Dataset (Iceberg)
        ↓  Funnel 增量索引 (批/流/CDC + 用户 Edit)
   Ontology Core: OMS(Schema) + OSv2(存储) + OSS(读) + Actions(写)
        ↓  OSDK / Ontology API / AIP Tool Call
   Applications: Workshop / Vertex / OSDK App  ↔  Human / Agent 决策
        ↓  Action 提交 (可 HITL 暂存)
   Writeback → Object DB 更新 + 审计日志 + 可选回写外部 ERP/MES

本质：人类与 AI Agent 在同一受治理的语义+执行平面上协作""")

    # --- Part 3 AIP ---
    s = new_content_slide(prs, "AIP 定位", "不是新大模型，是编排·治理·执行层")
    add_bullets(s, [
        "AIP = 架在 Foundry/Ontology 与 LLM 之间的企业 AI 操作系统",
        "核心命题：用 Ontology 给 LLM 划定「业务世界边界」",
        "         用确定性系统约束概率性模型输出",
        "五大模块：k-LLM · AIP Logic · Agent Studio · Assist · Evals",
        "全局锚定：所有模块读写必须经过 Ontology API（对象+链接+动作+函数）",
        "比喻：LLM = 可替换发动机；Ontology = 道路与交通规则；AIP = 方向盘+红绿灯+摄像头",
    ])

    s = new_content_slide(prs, "k-LLM 范式", "模型抽象 · 安全路由 · 多模型合议")
    add_bullets(s, [
        "两层含义（均属 AIP 公开表述）：",
        "① 模型抽象与路由层（Kernel + LLM）：",
        ("Conjure RPC 标准化接口，屏蔽 OpenAI/Anthropic/Llama 等 API 差异", 1),
        ("按任务复杂度、数据密级、成本延迟路由到合适模型", 1),
        ("敏感数据走本地私有模型；VPC/断网环境内推理，禁止数据回流训练", 1),
        "② K-LLM 多模型合议（AIPCon 2023 CTO 阐述）：",
        ("同一 Prompt 并行送入 K 个 LLM → Synthesis 阶段比较/评分/综合", 1),
        ("产出最优答案 + 异议视图；降低单模型幻觉与版本漂移风险", 1),
        ("即使 K=1，也可批量运行检测内部不一致性", 1),
        "推理视野受 Ontology 行列级权限（MAC/DAC）强制裁剪",
    ], font_size=14)

    s = new_content_slide(prs, "AIP Logic", "确定性 + 概率性混合编排")
    add_bullets(s, [
        "无代码 AI 函数构建器：Block 链式架构",
        "核心 Block：Use LLM · Execute Function · Apply Action · Conditionals · Loops",
        "关键设计：LLM 只「提议」工具调用参数，实际执行由 AIP 后端确定性完成",
        "Tools 绑定 Ontology：Query Objects / Call Function / Apply Action",
        "LLM 无法直接触碰 Dataset 或外部系统—必须经过权限校验",
        "Debugger：逐 Block 查看 CoT、工具请求、权限裁决结果",
        "Automate 集成：对象属性变更可触发 Logic（如供应链重排建议）",
    ])

    s = new_content_slide(prs, "AIP Agent Studio", "有界世界内的企业级 Agent")
    add_bullets(s, [
        "World Definition：静态声明可读 Object Types、可调用 Action/Function",
        "Agent「路径空间」预定义，杜绝越权探索（非松散 LangGraph 动态跳转）",
        "多 Agent 协作：共享 Ontology 状态 + Notification/Post-Function 触发",
        "Human-in-the-loop（HITL）强制插入：",
        ("Agent 写操作默认 Draft/Staged Action", 1),
        ("人类在 Workshop/Threads 审批后才提交 Ontology", 1),
        ("高风险 Action 可配置「必须人工确认」", 1),
        "发布：版本控制 + 粒度 RBAC + 使用监控 + 用户反馈（赞/踩）",
    ])

    s = new_content_slide(prs, "AIP Assist & AIP Evals")
    add_bullets(s, [
        "AIP Assist — 平台原生 Copilot：",
        ("上下文感知：知晓当前 App（Code Repo/Pipeline/Ontology Manager）", 1),
        ("编码辅助：Python/TS/SQL/PySpark + Palantir 特定 API", 1),
        ("继承 Foundry 安全模型，不访问无权限数据", 1),
        "AIP Evals — 非确定性输出的确定性测试：",
        ("测试用例集 + 评价指标（精确匹配/编辑距离/自定义评分）", 1),
        ("同测试集对比不同模型/Prompt/温度，量化方差", 1),
        ("指标达标才允许 Logic/Agent 发布生产", 1),
        ("+ AIP Observability：LLM输入→工具调用→权限裁决→输出全链路追踪", 1),
    ], font_size=15)

    s = new_content_slide(prs, "AIP 五模块 × Ontology 锚定")
    add_bullets(s, [
        "k-LLM：路由后的模型调用，读写必须经 Ontology API",
        "AIP Logic：Tools = Ontology Query / Function / Action",
        "Agent Studio：世界边界 = Object & Action Types 静态定义",
        "AIP Assist：辅助编写 Ontology/Foundry 代码与配置",
        "Evals：评估基于 Ontology 上下文的 LLM 输出准确性与安全性",
        "贯穿原则：Ontology-first，非 Data-first 或 Prompt-first",
    ])

    s = new_content_slide(prs, "完整交互时序图", "用户提问 → HITL 审批 → Action 写回")
    add_mono_block(s, """用户（Workshop / Threads）
  │  自然语言提问 / 委托任务
  ▼
AIP Assist / Agent 入口
  │  解析意图 · 加载 Agent World Definition（允许的 Object/Action 集）
  ▼
k-LLM 路由层
  │  按密级/复杂度选模型（或 K 模型并行）
  │  Ontology 权限裁剪 → 仅授权对象进入 Prompt 上下文
  ▼
AIP Logic（Block 链）
  │  Block1: Use LLM → 提议 Query Objects(parameters)
  │  Block2: Execute Function → OSS 查询 / 聚合 / Search Around
  │  Block3: Use LLM → 基于子图推理，提议 Apply Action(parameters)
  ▼
Agent Studio 提议层
  │  生成 Draft / Staged Action（默认不写生产）
  ▼
Human-in-the-Loop（Workshop / Threads）
  │  人工审阅：推理链 · 数据溯源 · Action 参数
  │  批准 / 驳回 / 修改
  ▼
Actions Service（唯一写入口）
  │  原子校验：参数 · 业务规则(Functions) · RBAC · 行列策略
  │  提交事务 → Funnel 索引 → OSv2 更新
  ▼
副作用与闭环
  ├─ 审计日志（Decision Lineage）
  ├─ Webhook 回写外部 ERP/MES（可选）
  ├─ Automate 触发下游 Logic/Agent
  └─ Evals/Observability 记录全链路""", size=10)

    # --- Part 4 Apollo ---
    s = new_content_slide(prs, "Apollo 重新定义", "自主部署操作系统 · 非传统 CI/CD")
    add_bullets(s, [
        "定位：基于控制论（Cybernetics）的自主部署操作系统",
        "≠ Jenkins / ArgoCD：不是执行固定 YAML 的 GitOps 流水线",
        "核心创新：约束求解（Constraint Solving）替代静态目标状态声明",
        "Hub-Spoke 异步架构：适应全球分布式、断网、气隙（Air-gapped）环境",
        "服务范围：Foundry / Gotham / AIP 等数百微服务的舰队级自治管理",
        "与 Rubix 关系：Apollo 编排层运行在 Rubix（加固 K8s 运行时）之上",
        "官方表述：「Mission Control for Software Deployment」",
    ])

    s = new_content_slide(prs, "Apollo Hub-Spoke 架构", "控制论模型")
    add_mono_block(s, """                    ┌──────────────── Apollo Hub（决策大脑）──────────────┐
                    │  Product Catalog · Release Channels · Orchestration │
                    │  Engine（约束求解）· Change Management（审批/审计）   │
                    └───────────────────────┬───────────────────────────┘
                                            │ Plan 下发（Hub 主动推送逻辑）
                    ┌───────────────────────┼───────────────────────────┐
                    │  Spoke 出站连接 ↑      │      ↑ Spoke 出站连接      │
        ┌───────────┴──────────┐  ┌────────┴────────┐  ┌────────────┴──────┐
        │ 公有云 K8s (Rubix)     │  │ 私有数据中心       │  │ 边缘/气隙/战术节点   │
        │ Spoke Control Plane    │  │ Spoke Agent       │  │ 离线 Artifact Bundle│
        │ Report State ↑         │  │ Pull & Execute Plan│  │ 签名验证后自主升级  │
        └────────────────────────┘  └───────────────────┘  └─────────────────────┘
关键：Spoke 主动出站连接 Hub — 无需边缘开放入站端口 · 天然兼容断网重连""")

    s = new_content_slide(prs, "约束求解编排引擎", "No Single Target State")
    add_bullets(s, [
        "Desired Intent：Product + Release Channel（如 Foundry :: Stable）",
        "  → 非固定版本号，而是「满足约束的最新可用版本」",
        "输入三要素（官方 Orchestration Engine 控制回路）：",
        ("① Product Catalog：可用版本 + Manifest 中的依赖/约束声明", 1),
        ("② Environment Settings：通道订阅、维护窗口、资源配额、审批策略", 1),
        ("③ Reported State：Spoke 回报的版本、健康度、Probe、Telemetry", 1),
        "硬约束：服务依赖图谱 · DB Schema 兼容性 · 安全 Markings",
        "软约束：维护窗口 · 带宽限制 · 业务高峰期禁止有损变更",
        "输出：Apollo Plan（仅当全部约束满足才下发执行）",
        "部分节点离线时：其余节点局部自治，重连后最终收敛",
    ], font_size=14)

    s = new_content_slide(prs, "Write Once · 气隙交付 · 自治运维")
    add_bullets(s, [
        "Product Specification & SDK：声明式规格（存储/IAM/网络），非手写 K8s YAML",
        "Service Management Plane：高层规格 → 翻译为 K8s CRD / VM 镜像 / 物理机包",
        "Rubix 基座：零信任 K8s · 多节点蓝绿部署 · 节点强制周期置换 · 全流量 mTLS",
        "气隙交付（官方支持多种模态）：",
        ("Cross Domain Solution (CDS) 跨域传输", 1),
        ("Delta 增量更新 + 边缘本地缓存", 1),
        ("加密签名离线 Artifact Bundle + 物理介质", 1),
        "Plan 生命周期：执行 → Probe 健康检查 → 失败自动 Rollback Plan",
        "Release Channel：Dev → Canary → Stable 自动推进（基于错误率等指标）",
        "合规：SAML/OIDC · RBAC 审批 · 不可变审计日志（FedRAMP/IL5-IL6 对齐）",
    ], font_size=14)

    s = new_content_slide(prs, "Apollo vs Kubernetes", "编排层 vs 基础设施层")
    add_bullets(s, [
        "K8s/Rubix 负责：容器生命周期、调度、网络、存储等基础设施原语",
        "Apollo 负责：产品级语义与舰队级治理",
        ("理解「Foundry = 300+ 微服务」的产品依赖图谱", 1),
        ("理解「Gotham 依赖特定 Ontology 版本」的跨产品约束", 1),
        ("理解「AIP 模型需 GPU 节点 + 出口管制合规」", 1),
        "ArgoCD/GitOps 擅长单集群固定 Manifest 收敛",
        "Apollo 擅长：跨成千上万异构集群/非 K8s 环境的复杂依赖与合规交付",
        "护城河：国防/情报领域「战火中断网维护大型分布式系统」经验产品化",
    ])

    s = new_content_slide(prs, "Foundry 微服务编排推演", "Apollo × Ontology × AIP 依赖对齐")
    add_mono_block(s, """【阶段 0】开发者发布
  各团队独立发布微服务 → 注册 Apollo Catalog（Manifest 声明依赖/Schema/Probe 约束）
  例：ontology-oss v3.2 依赖 ontology-osv2 >= 2.1；aip-logic 依赖 k-llm-router + GPU pool

【阶段 1】Hub 感知
  Catalog 新版本入库 → Orchestration Engine 读取各 Spoke Reported State
  Spoke-A（生产）：Foundry::Stable 通道 · 当前 200/300 服务已对齐 · OSS 已 v3.1

【阶段 2】约束求解（拓扑排序 + 约束剪枝）
  引擎构建产品依赖 DAG：
    Layer 0: Rubix substrate · Multipass/Alta 基础设施
    Layer 1: OMS · Object DB(OSv2) · Funnel · OSS（Ontology 核心读路径）
    Layer 2: Actions · Functions · Workshop backend
    Layer 3: AIP k-LLM Router · Logic Runtime · Agent Studio
    Layer 4: 面向用户的 Workshop / Vertex / AIP Threads
  检查：OSv2 迁移完成？Schema migration 窗口？Canary 错误率 < 阈值？
  → 暂停有损变更 或 生成有序 Plan 序列

【阶段 3】Plan 执行（Spoke Pull · Rubix 蓝绿）
  Plan-1: 升级 Funnel + OSv2（Ontology 写路径前置）
  Plan-2: 升级 OSS + OMS（读路径 · Schema 版本对齐）
  Plan-3: 升级 Actions（写入口 · 与 OSv2 兼容）
  Plan-4: 升级 AIP Logic + k-LLM（GPU 节点就绪后）
  每步：绿环境并行构建 → Probe 通过 → 流量渐进切换 → 销毁蓝节点

【阶段 4】收敛与通道推进
  全部 Probe 绿色 → Reported State 更新 → Stable 通道标记对齐
  Canary 环境 24h 错误率 < 0.1% → 自动 Promotion 到 Stable
  失败：Rollback Plan 回退 · 审计日志记录 · Change Management 通知

【与 Ontology/AIP 绑定】
  Ontology Schema 版本 = 产品约束（非独立部署可忽略）
  AIP Agent 仅在 Ontology + Logic 版本对齐后才可发布（Evals 门控）""", size=9.5)

    s = new_content_slide(prs, "全栈交付链路", "Apollo → Rubix → Foundry → Ontology → AIP")
    add_mono_block(s, """硬件/云基础设施
  ↓  Apollo Bootstrap（Hub/Spoke Environment 注册）
Rubix（零信任 K8s 运行时 · 蓝绿基座）
  ↓  Apollo Plan：Foundry Platform Services
Foundry Core（Pipeline · MMDP · Contour · Quiver · Workshop infra）
  ↓  Apollo Plan：Ontology Backend（OMS · Funnel · OSv2 · OSS · Actions）
Ontology Semantic Kernel（业务对象世界 · AI Grounding 边界）
  ↓  Apollo Plan：AIP Layer（k-LLM · Logic · Agent Studio · Evals）
AIP 编排治理层（LLM 提议 · 系统执行 · HITL 审批）
  ↓  用户/Agent 通过 Workshop · OSDK · Threads 操作
企业决策闭环（Action 写回 · Decision Lineage · 审计）

一句话：Apollo 保证「底层版本对齐」；Ontology 保证「语义与权限对齐」；AIP 保证「AI 行为对齐」""", size=11)

    s = new_content_slide(prs, "技术护城河总结")
    add_bullets(s, [
        "Foundry：数据操作系统—集成、治理、分析、运营应用一体化",
        "Ontology：可执行语义内核—Semantic+Kinetic+Governance 三区合一",
        "OSv2：索引/查询解耦 + Funnel + Spark 级 Search Around 规模化",
        "AIP：Ontology-first AI—LLM 提议、系统执行、人类审批、Evals 门控",
        "Apollo：约束求解舰队管理—Write Once Deploy Anywhere + 气隙 SaaS",
        "与纯 RAG/Agent/GitOps 差异：",
        ("持久对象世界 + 受控写入口 + 权限内生 + 全链路审计", 1),
        ("产品级依赖编排 + 断网自治 + 合规感知变更", 1),
        "一句概括：Palantir 卖的不是模型，是「企业决策操作系统 + 自主交付控制面」",
    ])

    s = new_content_slide(prs, "参考来源与说明")
    add_bullets(s, [
        "Palantir 官方文档：",
        ("foundry/object-backend · ontologies · vertex · aip · apollo/core · rubix", 1),
        "公开资料：Apollo White Paper · Palantir Blog（Apollo Orchestration）",
        ("AIPCon 2023（k-LLM）· PFCS Forward（气隙交付）· NVIDIA AIOS-RA", 1),
        "说明：",
        ("本文档为技术调研整理，非 Palantir 官方材料", 1),
        ("微服务编排推演链为基于公开架构的逻辑推演，非内部实现披露", 1),
        ("产品能力以 Palantir 当前文档为准", 1),
        "关联内部文档：M7-4 BDNS 生物数字自主进化智能体方案",
        "补充文档：Palantir-Apollo-深度解析.md",
    ], font_size=14)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
