# -*- coding: utf-8 -*-
"""Palantir 深度解析 PPTX v2 — 结构图为主 + 文字注解。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"c:\work\projects\wchat\docs\ref\Palantir-Foundry-AIP-Ontology-Apollo-解析和优化.pptx"
OUT_TMP = r"c:\work\projects\wchat\docs\ref\Palantir-Foundry-AIP-Ontology-Apollo-解析和优化-v2.2.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
TEAL = RGBColor(0x1A, 0x7A, 0x6E)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x88, 0x88, 0x88)
PALE_BLUE = RGBColor(0xD6, 0xEA, 0xF5)
PALE_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
PALE_GREEN = RGBColor(0xD5, 0xF0, 0xDC)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
RED = RGBColor(0xC0, 0x39, 0x2B)

BODY = 16
TITLE = 26
SUBTITLE = 16


def bg(slide, c=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c


def header(slide, title, sub=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.5), Inches(0.55))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(TITLE); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        s = slide.shapes.add_textbox(Inches(0.45), Inches(0.62), Inches(8.5), Inches(0.35))
        sp = s.text_frame.paragraphs[0]; sp.text = sub; sp.font.size = Pt(SUBTITLE); sp.font.color.rgb = RGBColor(0xAA,0xCC,0xEE)


def notes(slide, lines, left=7.35, top=1.25, w=5.7, h=6.0, sz=BODY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xDD,0xDD,0xDD)
    tb = slide.shapes.add_textbox(Inches(left+0.12), Inches(top+0.1), Inches(w-0.24), Inches(h-0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = DARK; p.space_after = Pt(4)


def rbox(slide, l, t, w, h, text, fill, tc=WHITE, sz=BODY, bold=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = tc; p.alignment = PP_ALIGN.CENTER
    return sh


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = GRAY; c.line.width = Pt(1.5)
    return c


def slide(title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, title, sub); return s


def pyramid_layer(s, l, t, w, h, text, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(BODY); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER


def cover():
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
    t = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.0))
    p = t.text_frame.paragraphs[0]; p.text = "Palantir 深度技术解析"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    s2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.8))
    sp = s2.text_frame.paragraphs[0]
    sp.text = "Foundry · Ontology · AIP · Apollo\n从「为什么」到「你的 L1/L2/L3 架构」"; sp.font.size = Pt(20); sp.font.color.rgb = RGBColor(0xAA,0xCC,0xEE)
    s3 = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(11), Inches(0.4))
    s3.text_frame.paragraphs[0].text = "谛听技术调研 · 2026-07-12 · 看图为主 · 文字为注解"; s3.text_frame.paragraphs[0].font.size = Pt(12); s3.text_frame.paragraphs[0].font.color.rgb = GRAY


# fix cover - use s not slide
def cover_fixed():
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.0))
    p = t.text_frame.paragraphs[0]; p.text = "Palantir 深度技术解析"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    s2 = s.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.8))
    sp = s2.text_frame.paragraphs[0]
    sp.text = "Foundry · Ontology · AIP · Apollo\nL1稳 + L2双引擎并行 + L3交互界面"; sp.font.size = Pt(20); sp.font.color.rgb = RGBColor(0xAA,0xCC,0xEE)


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

cover_fixed()

# 1 学习地图
s = slide("学习地图", "别记「怎么做」，记「为什么」")
rbox(s, 0.4, 1.4, 2.2, 0.85, "杂乱数据\nExcel/PDF/IoT", RGBColor(0xBB,0xBB,0xBB), DARK)
arrow(s, 2.65, 1.82, 3.2, 1.82)
rbox(s, 3.25, 1.3, 2.4, 1.05, "Foundry\n数据操作系统", BLUE, WHITE, bold=True)
arrow(s, 5.7, 1.82, 6.25, 1.82)
rbox(s, 6.3, 1.3, 2.5, 1.05, "Ontology\n可执行语义内核", TEAL, WHITE, bold=True)
arrow(s, 8.85, 1.82, 0.4, 2.65)
rbox(s, 3.0, 2.55, 2.8, 0.95, "AIP\nAI 编排治理层", PURPLE, WHITE, bold=True)
arrow(s, 5.85, 3.02, 6.4, 3.02)
rbox(s, 6.45, 2.55, 2.2, 0.95, "Apollo\n自主交付 OS", ORANGE, WHITE, bold=True)
rbox(s, 0.4, 3.5, 7.0, 3.0, "", WHITE, DARK)
for i, (t, c, y) in enumerate([
    ("为什么 Foundry？\n数据孤岛无法决策", PALE_BLUE, 3.65),
    ("为什么 Ontology？\n表没有「动作」和「权限」", PALE_GREEN, 4.55),
    ("为什么 AIP？\nLLM 需要世界边界", RGBColor(0xE8,0xDF,0xF5), 5.45),
]):
    rbox(s, 0.55, y, 6.5, 0.85, t, c, DARK)
notes(s, ["核心收获", "话术：先建本体模型，再构决策流", "比「我能写代码」高级", "对应任何政府/药企/电商项目", "今晚产出：你自己的 L1/L2/L3 白皮书"])

# 2 四大产品关系图
s = slide("四大产品关系", "不是四个工具，是一条价值链")
rbox(s, 0.5, 1.5, 6.7, 1.15, "Foundry — 企业数据操作系统\n集成 · 治理 · 管道 · 分析应用", BLUE, WHITE, bold=True)
arrow(s, 3.85, 2.7, 3.85, 3.05)
rbox(s, 0.5, 3.1, 6.7, 1.15, "Ontology — 可执行语义内核\nObject · Link · Action · Function · Governance", TEAL, WHITE, bold=True)
arrow(s, 3.85, 4.3, 3.85, 4.65)
rbox(s, 0.5, 4.7, 3.2, 1.1, "AIP\nLLM 编排治理", PURPLE, WHITE, bold=True)
rbox(s, 4.0, 4.7, 3.2, 1.1, "Apollo\n舰队交付 OS", ORANGE, WHITE, bold=True)
arrow(s, 2.1, 5.85, 2.1, 6.2)
rbox(s, 0.5, 6.25, 6.7, 0.75, "Rubix 加固 K8s 运行时（Apollo 编排其上）", GRAY, WHITE)
notes(s, ["为什么这样分层？", "Foundry 解决「有没有数据」", "Ontology 解决「数据像不像业务」", "AIP 解决「AI 能不能行动」", "Apollo 解决「能不能到处跑」", "Gotham 同源，偏国防情报"])

# 3 核心逻辑：数据本体化
s = slide("核心逻辑：数据如何变成「可计算对象」", "甲方话术：本体化 → 决策流")
for i, t in enumerate(["Excel", "PDF", "ERP", "IoT", "日志"]):
    rbox(s, 0.4+i*1.15, 1.5, 1.0, 0.7, t, RGBColor(0xCC,0xCC,0xCC), DARK)
arrow(s, 3.2, 2.25, 3.2, 2.6)
rbox(s, 0.5, 2.65, 6.5, 0.9, "① Connectors + Pipeline → Dataset (Iceberg/Parquet)", BLUE, WHITE)
arrow(s, 3.25, 3.6, 3.25, 3.95)
rbox(s, 0.5, 4.0, 6.5, 0.9, "② Funnel 映射 → Object / Property / Link", TEAL, WHITE)
arrow(s, 3.25, 4.95, 3.25, 5.3)
rbox(s, 0.5, 5.35, 6.5, 0.9, "③ Action + Function → 可执行决策", PURPLE, WHITE)
arrow(s, 3.25, 6.3, 3.25, 6.55)
rbox(s, 0.5, 6.6, 6.5, 0.75, "④ Workshop / AIP Agent → 人类+AI 行动", ORANGE, WHITE)
notes(s, ["为什么不用纯 RAG？", "表/chunk 没有类型和动作", "Ontology = 带权限的业务对象", "spf50 概念 → 链接 → 三款SKU", "话术：先本体模型，再决策流"])

# 4 Ontology 四层栈
s = slide("Ontology 四层栈", "Executable Semantic Kernel")
layers = [
    ("L4 应用/AI", "Workshop · Vertex · AIP Agent", PURPLE),
    ("L3 核心引擎", "OMS · OSv2 · OSS · Actions", TEAL),
    ("L2 映射层", "Funnel · CDC 增量索引", BLUE),
    ("L1 物理数据", "Backing Dataset (Iceberg)", GRAY),
]
for i, (t, d, c) in enumerate(layers):
    y = 1.35 + i * 1.2
    rbox(s, 0.5, y, 6.7, 1.05, f"{t}\n{d}", c, WHITE, bold=(i == 0))
    if i < 3:
        arrow(s, 3.85, y+1.05, 3.85, y+1.2)
notes(s, ["为什么不直接查表？", "底层数据不动", "物化索引低延迟读", "应用不碰 Parquet", "Semantic/Kinetic/Governance 三区并列"])

# 5 三区并列
s = slide("Ontology 三区", "运行时原子裁决")
rbox(s, 0.5, 1.55, 2.1, 2.5, "Semantic\n名词\nObject\nLink\nInterface", TEAL, WHITE, bold=True)
rbox(s, 2.75, 1.55, 2.1, 2.5, "Kinetic\n动词\nAction\nFunction", ORANGE, WHITE, bold=True)
rbox(s, 5.0, 1.55, 2.1, 2.5, "Governance\n治理\nMarkings\nRBAC", NAVY, WHITE, bold=True)
rbox(s, 0.5, 4.25, 6.6, 0.85, "AI Agent 只能通过 Action/Function 触碰世界 — 不能直读 Dataset", PURPLE, WHITE)
notes(s, ["为什么三区并列？", "名词/动词/权限一次裁决", "避免 LLM 越权写库", "Decision Lineage 全审计", "这是 AIP 的根基"])

# 6 OSv2 vs Phonograph
s = slide("OSv2 架构演进", "索引与查询解耦")
rbox(s, 0.5, 1.5, 3.0, 2.8, "OSv1 Phonograph\n(遗留)\n索引+查询耦合\n2026 弃用", RGBColor(0xAA,0xAA,0xAA), WHITE)
rbox(s, 3.7, 1.5, 3.5, 2.8, "OSv2\n索引 ‖ 查询\nFunnel+CDC\nSpark Search Around", TEAL, WHITE, bold=True)
arrow(s, 3.55, 2.9, 3.7, 2.9)
notes(s, ["为什么重写？", "数十亿对象规模", ">100K Search Around → Spark", "增量索引全对象类型", "官方强制迁移"])

# 7 AIP 五模块
s = slide("AIP 五模块", "全部锚定 Ontology")
rbox(s, 2.8, 2.45, 2.8, 1.35, "Ontology\n对象+动作+函数", TEAL, WHITE, bold=True)
mods = [("k-LLM\n路由", 0.6, 1.55), ("Logic\n编排", 5.3, 1.55), ("Agent\nStudio", 0.6, 4.1), ("Assist", 5.3, 4.1), ("Evals\n门控", 2.8, 5.05)]
for t, x, y in mods:
    rbox(s, x, y, 1.9, 0.95, t, PURPLE, WHITE)
    arrow(s, x+0.95, y+0.95 if y>2 else y-0.05, 2.8 if x<3 else 4.2, 2.45 if y<2 else 3.55)
notes(s, ["为什么不是 ChatGPT 套壳？", "LLM 只提议，系统执行", "HITL 默认暂存 Action", "Evals 达标才上生产", "k-LLM = 路由+多模型合议"])

# 8 Apollo Hub-Spoke
s = slide("Apollo Hub-Spoke", "控制论 · 非 Jenkins")
rbox(s, 2.3, 1.4, 3.6, 1.15, "Hub\nCatalog · 约束求解引擎", ORANGE, WHITE, bold=True)
for i, t in enumerate(["公有云", "私有云", "气隙边缘"]):
    rbox(s, 0.5+i*2.35, 3.1, 2.1, 1.2, f"Spoke\n{t}\n出站Pull", PALE_ORANGE, DARK)
    arrow(s, 2.5+i*2.35+1.05, 3.05, 4.1, 2.6)
rbox(s, 0.5, 4.55, 6.6, 0.8, "无单一目标状态 — Product+Channel+约束 → Plan", NAVY, WHITE)
notes(s, ["为什么 Spoke 出站？", "边缘无需开入站端口", "断网重连后收敛", "气隙：签名 Artifact Bundle", "≠ K8s 替代品，是上层编排"])

# 9 全栈链路
s = slide("全栈交付链路")
steps = ["硬件", "Apollo\nBootstrap", "Rubix", "Foundry", "Ontology", "AIP", "用户/Agent"]
for i, t in enumerate(steps):
    rbox(s, 0.3+i*1.05, 2.0, 1.0, 1.35, t, [BLUE,ORANGE,GRAY,BLUE,TEAL,PURPLE,GREEN][i], WHITE)
    if i < 6:
        arrow(s, 0.3+i*1.05+1.0, 2.67, 0.3+(i+1)*1.05, 2.67)
notes(s, ["Apollo 保证版本对齐", "Ontology 保证语义对齐", "AIP 保证行为对齐", "300+ 微服务 DAG 拓扑 Plan", "失败自动 Rollback"])

# 10 站在巨人肩膀上
s = slide("站在巨人的肩膀上", "高性能实时决策架构 · 双引擎并行")
rbox(s, 0.5, 1.4, 6.7, 1.05, "L1 离线本体层\n对标 Palantir Ontology\n直接复用 · 记性好 · 知识底座", BLUE, WHITE, bold=True)
rbox(s, 0.5, 2.6, 6.7, 1.05, "L3 交互界面\n对标 Palantir Apollo\n直接复用 · 数字孪生大屏 · 运营后台", ORANGE, WHITE, bold=True)
rbox(s, 0.5, 3.8, 6.7, 1.45, "L2 决策层（双引擎并行）\n左：Palantir AIP(K-LLM) 主路  右：自研引擎(PPR+ANN+EGB) 应急车道\n不是替代，是补充 · 互为备份", TEAL, WHITE, bold=True)
rbox(s, 0.5, 5.45, 6.7, 0.9, "金句：不重复造轮子，只造更强的引擎", GOLD, DARK, bold=True)
notes(s, ["L1+L3 复用最佳实践", "L2 双引擎并列", "AIP=主路 自研=应急", "稳+快 双保险", "M7-6 v2.0"])

# 11 架构金字塔
s = slide("架构金字塔", "L1 稳 · L2 双引擎 · L3 交互界面")
pyramid_layer(s, 1.0, 1.45, 5.8, 1.0, "L3 交互界面\n对标 Apollo · 可视化交付", ORANGE)
pyramid_layer(s, 0.5, 2.6, 6.8, 1.35, "L2 决策层（双引擎）\nAIP主路 ‖ 自研应急车道", TEAL)
pyramid_layer(s, 0.2, 4.1, 7.4, 1.25, "L1 离线本体层\n对标 Ontology · 知识底座", BLUE)
rbox(s, 0.5, 5.55, 6.7, 0.8, "Palantir 负责「稳」，自研引擎负责「快」", NAVY, WHITE, bold=True)
notes(s, ["塔基 L1 本体", "塔身 L2 双引擎", "塔尖 L3 交互界面", "并列不替代", "稳快结合"])

# 12 痛点：水土不服
s = slide("痛点：AIP 主路在电商洪峰的压力", "需要 L2 右引擎应急补位，不是替换 AIP")
pains = [
    ("高并发", "直播间一秒几千弹幕\nAIP 检索链路太长 → 卡顿", RED),
    ("新场景", "「绝绝子」「City不City」\n规则库更新跟不上", ORANGE),
    ("成本高", "全量走 AIP 重型接口\n财务难以承受", PURPLE),
]
for i, (t, d, c) in enumerate(pains):
    rbox(s, 0.5, 1.4 + i * 1.55, 6.7, 1.35, f"{t}\n{d}", c, WHITE, bold=True)
rbox(s, 0.5, 6.1, 6.7, 0.85, "结论：L2 并列加速引擎，分流减压", GREEN, WHITE, bold=True)
notes(s, ["AIP 仍是主路", "洪峰切右引擎", "99% 噪音先过滤", "各走各的链路", "M7-6 §4"])

# 13 三层架构详解（重点 · 双引擎并列图）
s = slide("方案：三层架构详解", "★ 重点页 · L2 双引擎并列")
rbox(s, 0.5, 1.3, 6.7, 0.95, "L3 交互界面（对标 Palantir Apollo）\n数字孪生大屏 · 运营操作后台", ORANGE, WHITE, bold=True)
rbox(s, 0.5, 2.35, 6.7, 0.55, "L2 决策层（双引擎并行）— 两个模块并列，互不冲突", TEAL, WHITE, bold=True)
rbox(s, 0.5, 3.0, 3.2, 1.55, "Palantir 原生 AIP\n(K-LLM)\n通用场景标准化决策\n检索→生成 【主路】", PURPLE, WHITE, bold=True)
rbox(s, 4.0, 3.0, 3.2, 1.55, "自研实时决策引擎\n(PPR+ANN+EGB)\n高并发场景加速\n过滤→提纯→快决策【应急车道】", GREEN, WHITE, bold=True)
arrow(s, 3.75, 3.77, 4.0, 3.77)
rbox(s, 0.5, 4.75, 6.7, 1.05, "L1 离线本体层（对标 Palantir Ontology）\nOKF · 知识底座 · 老板熟悉的「知识底座」", BLUE, WHITE, bold=True)
notes(s, ["左引擎=AIP主路", "右引擎=应急补位", "环科院→走AIP", "马帮直播→走右引擎", "互为备份"])

# 14 L2 双引擎运作逻辑
s = slide("L2 双引擎运作逻辑", "主路 + 应急车道 · 各走各的")
rbox(s, 0.5, 1.35, 3.2, 0.7, "左引擎 AIP 主路", PURPLE, WHITE, bold=True)
rbox(s, 4.0, 1.35, 3.2, 0.7, "右引擎 自研应急车道", GREEN, WHITE, bold=True)
rbox(s, 0.5, 2.15, 3.2, 3.5, "① Ontology 检索\n② K-LLM 生成\n③ HITL 审批写回\n\n场景：\n政策查询\n常规客服咨询", PURPLE, WHITE)
rbox(s, 4.0, 2.15, 3.2, 3.5, "① PPR 并行分流\n② ANN 特征萃取\n③ EGB 动态调优\n\n场景：\n直播间弹幕\n订单洪峰", GREEN, WHITE)
rbox(s, 0.5, 5.85, 6.7, 0.85, "平时 AIP 主力 · 洪峰自动切右引擎 · 双保险", NAVY, WHITE, bold=True)
notes(s, ["检索→生成", "过滤→提纯→快决策", "交通指挥官PPR", "读心神探ANN", "进化引擎EGB"])

# 15 术语去术语化
s = slide("术语去术语化", "L2 双引擎对照表")
pairs = [
    ("AIP K-LLM", "主路决策引擎", "L2左"),
    ("PPR 并行模式识别", "交通指挥官", "L2右"),
    ("ANN 特征萃取引擎", "读心神探", "L2右"),
    ("EGB 动态调优机制", "进化引擎", "L2右"),
    ("Apollo/图谱", "交互界面", "L3"),
]
for i, (a, b, l) in enumerate(pairs):
    y = 1.35 + i * 0.95
    rbox(s, 0.5, y, 2.5, 0.85, a, PURPLE if "L2左" in l else TEAL, WHITE)
    rbox(s, 3.15, y, 2.5, 0.85, b, BLUE, WHITE)
    rbox(s, 5.8, y, 1.2, 0.85, l, ORANGE, WHITE)
notes(s, ["不是替代AIP", "是并列补充", "对外用中文名", "对内见 M7-4", "稳快双保险"])

# 16 零售案例
s = slide("零售/电商案例", "预制菜原料 · 注意公开来源边界")
cases = [
    ("C&A 时尚零售", "3月上线\n库存数字孪生\n+200M销售管理"),
    ("匿名时尚零售商", "巴西\n+16M价值\nSKU可用+13%"),
    ("F&B 快消", "减过库存\n实时运营画面"),
    ("跨境电商(你)", "马帮方案\nListing+库存Buddy"),
]
for i, (t, d) in enumerate(cases):
    x = 0.4 + (i%2)*3.5; y = 1.35 + (i//2)*2.1
    rbox(s, x, y, 3.3, 1.95, f"{t}\n{d}", [PALE_BLUE,PALE_GREEN,PALE_ORANGE,RGBColor(0xE8,0xDF,0xF5)][i], DARK)
notes(s, ["沃尔玛：自有 Retail Link", "非公开 Palantir 客户", "可口可乐：无公开案例", "可引用：C&A PR", "痛点：库存周转/定价/画像", "→ 马帮跨境电商预制菜"])

# 17 定价交付
s = slide("定价与交付模式", "Land-and-Expand")
rbox(s, 0.5, 1.45, 2.1, 1.65, "Layer1\n平台年费\n$5-25M/年", BLUE, WHITE, bold=True)
rbox(s, 2.75, 1.45, 2.1, 1.65, "Layer2\n扩展消费\n数据源/用户\nOntology/AIP", TEAL, WHITE, bold=True)
rbox(s, 5.0, 1.45, 2.1, 1.65, "Layer3\n专业服务\nFDE 驻场\n$200-500K/人", ORANGE, WHITE, bold=True)
rbox(s, 0.5, 3.25, 6.6, 1.05, "Pilot $1-5M → Enterprise $50-100M+/年 · NRR>118%", NAVY, WHITE)
rbox(s, 0.5, 4.5, 6.6, 2.0, "你的接项目建议\n年费(平台) + 实施费(本体建模) + 按SKU/店铺扩展\n先 2周 POC → 3月 MVP → 年框", GREEN, WHITE)
notes(s, ["Palantir 卖「决策能力」", "不是按席位卖 SaaS", "FDE 前向部署工程师", "高切换成本", "你可：本体建模费+年维护"])

# 18 时序图
s = slide("AIP × Ontology 时序", "用户提问 → Action 写回")
flow = ["用户", "k-LLM", "Logic", "Agent", "HITL", "Action", "反馈"]
for i, t in enumerate(flow):
    rbox(s, 0.3+i*1.0, 2.15, 0.95, 0.85, t, PURPLE if i<5 else TEAL, WHITE)
    if i < 6:
        arrow(s, 0.3+i*1.0+0.95, 2.57, 0.3+(i+1)*1.0, 2.57)
rbox(s, 0.5, 3.25, 6.6, 2.8, "LLM 提议 → 系统确定性执行 → 人类审批 → Funnel 索引 → 边权演化", PALE_GREEN, DARK)
notes(s, ["完整闭环", "权限裁剪后的对象进 Prompt", "Draft Action 默认", "Decision Lineage", "对应导购/文案 Buddy"])

# 19 价值总结
s = slide("价值总结", "双引擎并行 · 稳快双保险")
rbox(s, 0.5, 1.4, 2.1, 2.2, "对外\nPalantir 故事\nL1本体+L3交互界面\n全球最佳实践", BLUE, WHITE, bold=True)
rbox(s, 2.8, 1.4, 2.1, 2.2, "对内\nL2 双引擎\nAIP主路+\n自研应急补位", TEAL, WHITE, bold=True)
rbox(s, 5.1, 1.4, 2.1, 2.2, "结果\n高并发实时\n决策系统\n稳快双保险", GREEN, WHITE, bold=True)
rbox(s, 0.5, 3.85, 6.7, 2.6, "不是替代 AIP，是并列补充\n\nAIP 走「检索→生成」主路（政策查询·常规客服）\n自研引擎走「过滤→提纯→快决策」应急车道（直播弹幕·洪峰）\n\nPalantir 负责「稳」，自研引擎负责「快」\n平时 AIP 主力，洪峰自动补位，双保险", NAVY, WHITE, bold=True)
notes(s, ["M7-5 白皮书", "M7-6 v2.0 双引擎", "M7-4 BDNS", "本 PPT v2.2", "解析和优化版"])

for path in (OUT, OUT_TMP):
    try:
        prs.save(path)
        print("Saved:", path, "slides:", len(prs.slides))
        break
    except PermissionError:
        print("Locked, skip:", path)
else:
    raise SystemExit("无法写入 PPT，请关闭已打开的 Palantir PPTX 后重跑脚本")
