# -*- coding: utf-8 -*-
"""高性能实时决策架构 PPTX — 6 页对外宣讲 · 结构图 + 16pt 正文。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"c:\work\projects\wchat\docs\ref\高性能实时决策架构-站在巨人肩膀上.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
TEAL = RGBColor(0x1A, 0x7A, 0x6E)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x88, 0x88, 0x88)
PALE_BLUE = RGBColor(0xD6, 0xEA, 0xF5)
PALE_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)
PALE_GREEN = RGBColor(0xD5, 0xF0, 0xDC)
RED = RGBColor(0xC0, 0x39, 0x2B)

BODY = 16
TITLE = 26
SUBTITLE = 16


def bg(slide, c=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c


def header(slide, title, sub=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(12), Inches(0.55))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(TITLE); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        s = slide.shapes.add_textbox(Inches(0.45), Inches(0.62), Inches(12), Inches(0.35))
        sp = s.text_frame.paragraphs[0]; sp.text = sub; sp.font.size = Pt(SUBTITLE); sp.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)


def rbox(slide, l, t, w, h, text, fill, tc=WHITE, bold=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(BODY); p.font.bold = bold; p.font.color.rgb = tc; p.alignment = PP_ALIGN.CENTER
    return sh


def arrow(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = GRAY; c.line.width = Pt(2)
    return c


def slide(title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); header(s, title, sub); return s


def pyramid_layer(s, l, t, w, h, text, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(BODY); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 封面
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
t = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(1.2))
p = t.text_frame.paragraphs[0]
p.text = "高性能实时决策架构"; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
s2 = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.0))
sp = s2.text_frame.paragraphs[0]
sp.text = "基于 Palantir 思想 · 站在巨人的肩膀上"; sp.font.size = Pt(22); sp.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
rbox(s, 1.5, 4.2, 10.3, 1.2, "Palantir L1（本体库）+ L3（数字孪生）= 全球最佳实践，直接复用\n自研 L2（中国芯加速引擎）= 我们不重复造轮子，我们只造更强的引擎", GOLD, DARK, bold=True)
s3 = s.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11), Inches(0.4))
s3.text_frame.paragraphs[0].text = "谛听技术 · M7-6 · 2026-07-13"; s3.text_frame.paragraphs[0].font.size = Pt(BODY); s3.text_frame.paragraphs[0].font.color.rgb = GRAY

# 第2页：金字塔
s = slide("架构金字塔", "L1 稳 · L2 快 · L3 帅")
pyramid_layer(s, 3.8, 1.5, 5.5, 1.15, "L3 数字孪生层（交互）\n可视化 · 图谱 · 人机审批", ORANGE)
pyramid_layer(s, 2.6, 2.85, 7.9, 1.25, "L2 在线决策层（引擎）\n并行模式识别 · 特征萃取 · 动态调优", TEAL)
pyramid_layer(s, 1.2, 4.3, 10.7, 1.45, "L1 离线本体层（底盘）\nOKF + Palantir Ontology 思想 · 记性好 · 沉淀规则", BLUE)
rbox(s, 0.5, 6.0, 12.3, 0.85, "金句：Palantir 负责「稳」，咱们负责「快」—— 稳快结合，就是核心竞争力", NAVY, WHITE, bold=True)

# 第3页：痛点
s = slide("痛点：纯 Palantir 在国内电商的「水土不服」", "结论：L1 与 L3 之间，必须加「实时过滤与加速层」")
pains = [
    ("高并发", "直播间每秒上万弹幕\n检索链路太长 → 卡顿", RED, WHITE),
    ("新场景", "「绝绝子」「City不City」\n规则库更新跟不上", ORANGE, WHITE),
    ("成本高", "全量走重型 AI 接口\n财务难以承受", PURPLE, WHITE),
]
for i, (t, d, c, tc) in enumerate(pains):
    rbox(s, 0.5 + i * 4.2, 1.5, 3.9, 2.2, f"{t}\n\n{d}", c, tc, bold=True)
rbox(s, 0.5, 4.1, 12.3, 1.0, "纯 Palantir 擅长「世界模型」与「企业级治理」，但不擅长「电商洪峰下的毫秒分流」", NAVY, WHITE)
rbox(s, 0.5, 5.35, 12.3, 1.35, "我们需要在 L1 和 L3 之间，加一个 L2「实时过滤与加速层」\n让 99% 噪音在到达重型 AI 之前被拦截 · 分流 · 粗判", GREEN, WHITE, bold=True)

# 第4页：三层架构（重点）
s = slide("方案：三层架构详解", "★ 重点页")
rbox(s, 0.5, 1.4, 12.3, 1.45, "L1 离线本体层（底盘）\n你的库 + Palantir Ontology 思想\n负责「记性好」—— 沉淀知识、规则、对象关系", BLUE, WHITE, bold=True)
arrow(s, 6.65, 2.9, 6.65, 3.15)
rbox(s, 0.5, 3.2, 12.3, 1.75, "L2 在线决策层（引擎）★ 自研「中国芯」加速中间件\n并行模式识别（PPR）· 特征萃取引擎（ANN）· 动态调优机制（EGB）\n负责「反应快」—— 毫秒级分流、语义理解、策略自优化", TEAL, WHITE, bold=True)
arrow(s, 6.65, 5.0, 6.65, 5.25)
rbox(s, 0.5, 5.3, 12.3, 1.45, "L3 数字孪生层（交互）\n你的可视化 — 图谱 Tab · WorkBuddy · 运营看板\n负责「长得帅」—— 直观展示结果、溯源路径、人机审批", ORANGE, WHITE, bold=True)

# 第5页：L2 运作逻辑
s = slide("L2 层运作逻辑", "给 Palantir 装一套「中国芯」")
rbox(s, 0.5, 1.35, 12.3, 0.85, "输入：弹幕 · 订单流 · 库存数据（洪水般涌入）", RGBColor(0xBB, 0xBB, 0xBB), DARK, bold=True)
arrow(s, 6.65, 2.25, 6.65, 2.5)
steps = [
    ("① 并行模式识别（PPR）", "多车道分流\n「下单」和「闲聊」瞬间分开\n交通指挥官 · 绝不堵车", TEAL),
    ("② 特征萃取引擎（ANN）", "emoji · 错别字 · 方言 · 新梗\n抓出「愤怒」「想退货」「急件」\n读心神探 · 模糊→精确", PURPLE),
    ("③ 动态调优机制（EGB）", "记录每次行动结果\n淘汰差策略 · 保留好策略\n进化引擎 · 越跑越准", GREEN),
]
for i, (t, d, c) in enumerate(steps):
    y = 2.55 + i * 1.45
    rbox(s, 0.5, y, 12.3, 1.3, f"{t}\n{d}", c, WHITE, bold=(i == 0))
    if i < 2:
        arrow(s, 6.65, y + 1.3, 6.65, y + 1.45)
rbox(s, 0.5, 6.55, 12.3, 0.75, "输出：明确指令 — 回复 · 报警 · 改价 · 转人工 · 写回 L3", ORANGE, WHITE, bold=True)

# 第6页：价值总结
s = slide("价值总结", "对外高大上 · 对内低成本高速度")
rbox(s, 0.5, 1.4, 3.9, 2.5, "对外\n讲 Palantir 的故事\nL1 本体 + L3 数字孪生\n全球最佳实践", BLUE, WHITE, bold=True)
rbox(s, 4.7, 1.4, 3.9, 2.5, "对内\n自研 L2 加速引擎\nPPR + ANN + EGB\n低成本 · 高速度", TEAL, WHITE, bold=True)
rbox(s, 8.9, 1.4, 3.9, 2.5, "结果\n全球领先的\n高并发实时\n决策系统", GREEN, WHITE, bold=True)
rbox(s, 0.5, 4.2, 12.3, 2.5, "咱们的 L2 就是给 Palantir 装了一套「中国芯」\n\n靠并行模式识别（PPR）抗住电商高流量\n靠特征萃取引擎（ANN）看懂客户复杂需求\n靠动态调优机制（EGB）让系统越跑越快\n\nPalantir 负责「稳」，咱们负责「快」—— 稳快结合，就是咱们的核心竞争力", NAVY, WHITE, bold=True)

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides))
