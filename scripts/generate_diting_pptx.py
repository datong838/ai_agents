#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Diting bio-digital knowledge network PPTX from M7-2 script."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG = RGBColor(0x0A, 0x16, 0x28)
BG_CARD = RGBColor(0x12, 0x24, 0x40)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xD6, 0xE0)
MUTED = RGBColor(0x8A, 0x9B, 0xAD)
GREEN = RGBColor(0x06, 0xD6, 0xA0)
ORANGE = RGBColor(0xFF, 0x9F, 0x1C)
APPENDIX_BG = RGBColor(0x0D, 0x1A, 0x2E)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
TABLE_HEADER_BG = RGBColor(0x00, 0x6D, 0x87)
TABLE_ROW_BG = RGBColor(0xF7, 0xFA, 0xFC)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF2, 0xF7)
FONT = "Microsoft YaHei"
FONT_MIN = 16  # 正文/表格最小字号


def font_pt(legacy: int) -> Pt:
    """设计稿 legacy 字号 → 实际渲染 pt（正文不低于 FONT_MIN，标题顺延放大）。"""
    scale = {
        9: 16, 10: 16, 11: 16,
        12: 18, 13: 18,
        14: 20, 15: 20,
        16: 22, 17: 22, 18: 24,
        20: 26, 22: 28, 24: 30,
        28: 36, 32: 40, 36: 44, 38: 46, 40: 48,
    }
    return Pt(scale.get(legacy, max(legacy + 4, FONT_MIN)))

OUTPUT = (
    Path(__file__).resolve().parent.parent
    / "ref"
    / "生物-数字双模态知识网机制的智能体研发与应用.pptx"
)

MAIN_COUNT = 26
APPENDIX_START = 27


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def txt(slide, l, t, w, h, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = font_pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def bullets(slide, l, t, w, h, lines, size=13, color=LIGHT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = font_pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def line(slide, top=Inches(1.55)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(1.1), Inches(0.035))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def title(slide, text, sub=None):
    txt(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7), text, 36, WHITE, True)
    line(slide)
    if sub:
        txt(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.4), sub, 18, ACCENT2)


def footer(slide, num, tag="宣讲"):
    txt(slide, Inches(0.6), Inches(7.05), Inches(4), Inches(0.3), f"谛听 DITING · {tag}", 11, MUTED)
    txt(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.3), str(num), 11, MUTED, align=PP_ALIGN.RIGHT)


def _style_table_cell(cell, text, *, size, color, bold=False, bg):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = font_pt(size)
        p.font.color.rgb = color
        p.font.bold = bold


def table(slide, l, t, w, h, headers, rows, col_w=None):
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, l, t, w, h).table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for j, htxt in enumerate(headers):
        _style_table_cell(
            tbl.cell(0, j), htxt, size=14, color=WHITE, bold=True, bg=TABLE_HEADER_BG
        )
    for i, row in enumerate(rows, 1):
        row_bg = TABLE_ROW_BG if i % 2 == 1 else TABLE_ROW_ALT
        for j, cell in enumerate(row):
            _style_table_cell(
                tbl.cell(i, j), cell, size=11, color=BLACK, bold=False, bg=row_bg
            )
    return tbl


def card(slide, l, t, w, h, heading, body, hcolor=ACCENT):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = BG_CARD
    r.line.color.rgb = ACCENT
    txt(slide, l + Inches(0.15), t + Inches(0.12), w - Inches(0.3), Inches(0.35), heading, 16, hcolor, True)
    txt(slide, l + Inches(0.15), t + Inches(0.45), w - Inches(0.3), h - Inches(0.55), body, 13, LIGHT)


# ── Main slides ───────────────────────────────────────────

def s01_cover(prs):
    s = blank(prs)
    bg(s)
    txt(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5), "谛听  DITING", 15, ACCENT, True)
    txt(s, Inches(0.6), Inches(1.9), Inches(12), Inches(1.3), "让每一条知识都被「看见」", 40, WHITE, True)
    txt(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.9),
        "基于生物-数字双模态知识网与神经扩散机制的智能体系统研发及应用", 18, ACCENT2)
    txt(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.45),
        "生物-数字双模态知识网机制的智能体研发与应用", 14, MUTED)
    txt(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.4), "2026 · 谛听科技", 12, MUTED)
    footer(s, 1)


def s02_dark_age(prs):
    s = blank(prs)
    bg(s)
    title(s, "数据的黑盒时代", "企业正在被数据淹没，而不是被数据赋能")
    stats = [("70%+", "企业数据从未用于分析"), ("↑", "数据量与系统数持续增长"), ("隐性损失", "数据质量差 > 缺数据")]
    for i, (n, d) in enumerate(stats):
        x = Inches(0.6 + i * 4.1)
        txt(s, x, Inches(2.0), Inches(3.5), Inches(0.7), n, 36, ACCENT, True)
        txt(s, x, Inches(2.8), Inches(3.5), Inches(0.8), d, 13, LIGHT)
    bullets(s, Inches(0.6), Inches(4.0), Inches(12), Inches(2.5), [
        "工艺卡片 · 实验记录 · 商品详情 · 客服对话 · 专利文献 · 培训材料……",
        "问题不是「有没有数据」，而是「能否转化为行动」",
        "知识在专家脑中、PDF 里、各系统记录里 — 组织层面看不见、连不起来、用不了",
    ], 15)
    footer(s, 2)


def s03_dilemma(prs):
    s = blank(prs)
    bg(s)
    title(s, "传统方案的困局", "CRM / ERP / BI / 数据湖 · 数据孤岛 · AI 停在 Chat")
    boxes = [
        ("数据孤岛", "N 个系统需 N×(N-1)/2 连接\n复杂度指数级增长"),
        ("功能堆叠 ≠ 价值", "SaaS 卖 Feature 不卖决策\n10 个工具无法形成统一视图"),
        ("AI 停在聊天", "大模型触达不了真实业务对象\n无法触发业务动作"),
    ]
    for i, (h, b) in enumerate(boxes):
        card(s, Inches(0.6 + i * 4.1), Inches(1.8), Inches(3.9), Inches(2.2), h, b)
    txt(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
        "结构性困局：数据有了，语义没有；模型有了，本体没有；对话有了，行动没有", 15, ORANGE, True)
    footer(s, 3)


def s04_direction(prs):
    s = blank(prs)
    bg(s)
    title(s, "我们的方向", "智能体操作系统 · 算法渊源 + 工程落地")
    items = [
        ("生物-数字双模态", "节点=神经元 · 边=突触 · 查询=动作电位 · 反馈=突触可塑性"),
        ("知识网", "document/section/concept/media 异构图\n有向、有权、可解释的多层联结"),
        ("神经扩散机制", "ANN 种子 → meta-path + PPR 扩散\n经验传导，非黑盒检索"),
    ]
    for i, (h, b) in enumerate(items):
        y = Inches(1.55 + i * 1.35)
        txt(s, Inches(0.6), y, Inches(2.8), Inches(0.35), h, 15, ACCENT, True)
        txt(s, Inches(3.5), y, Inches(9), Inches(1.0), b, 12, LIGHT)
    card(s, Inches(0.6), Inches(5.35), Inches(12), Inches(1.35), "算法渊源",
         "研究起点：异构图神经网络(HIN/GNN) + 知识图谱 Message Passing\n"
         "工程现状：BGE + PPR + meta-path = 无需训练的 Message Passing（可解释·可部署）\n"
         "演进 L3：R-GCN / HAN 可学习扩散 · 共用 propagate 接口 · Skill 层不感知")
    txt(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
        "已落地：谛听客户端 + 营销智能体 · 需求雷达 / 知识图谱 / 客服智能体", 11, ACCENT2)
    footer(s, 4)


def s05_philosophy(prs):
    s = blank(prs)
    bg(s)
    title(s, "核心理念：增强人类，而非替代")
    pillars = [
        ("知识网本体", "document/section/concept/media 映射真实世界\n不是数据库表，是活的业务语义模型"),
        ("Human-in-the-Loop", "AI 分析推荐 · 人类判断决策\n战斗机 HUD，非自动驾驶"),
        ("Data → Decision → Action", "从「我知道了」到「我做了」\n话术 / 商品 / 方案 / 排查清单"),
    ]
    for i, (h, b) in enumerate(pillars):
        x = Inches(0.6 + i * 4.1)
        card(s, x, Inches(1.7), Inches(3.9), Inches(2.8), h, b)
    footer(s, 5)


def s06_compare(prs):
    s = blank(prs)
    bg(s)
    title(s, "谛听 ≠ 传统方案")
    table(s, Inches(0.6), Inches(1.6), Inches(12), Inches(4.5),
          ["维度", "传统方案", "谛听"],
          [
              ["核心逻辑", "功能模块堆叠", "知识网本体"],
              ["数据处理", "ETL + 数仓", "Pipeline + 实体抽取 + 图谱"],
              ["AI 策略", "独立 Chatbot", "AI 嵌入业务流 + 决策链"],
              ["集成", "API 点对点", "统一知识网 + 自动挂接"],
              ["扩展", "买新模块", "复用本体，边际成本趋零"],
              ["价值", "看报表", "做决策、做行动"],
          ],
          [Inches(1.8), Inches(3.5), Inches(6.7)])
    footer(s, 6)


def s07_layers(prs):
    s = blank(prs)
    bg(s)
    title(s, "三大能力层", "Knowledge Network · Neural Diffusion · Domain Agents")
    layers = [
        ("知识网层 · 学", "关系型数据/文档/图片/对话 → 异构图谱\nBGE-M3 向量 · VLM 描述 · contains/mentions/appears_in"),
        ("传导引擎层 · 想", "Query → Embedding → ANN → Meta-path + PPR → 子图\n向量相似 + 图结构传导 · 可解释可审计"),
        ("智能体层 · 做", "需求雷达 · 智能客服 · 内容写作 · 认知推演\nactivate → propagate → package → llm → skill_output"),
    ]
    for i, (h, b) in enumerate(layers):
        x = Inches(0.6 + i * 4.1)
        card(s, x, Inches(1.7), Inches(3.9), Inches(3.5), h, b)
    txt(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
        "共享底座：异构图 · 向量库 · LLM 网关 · 私有化部署 · 数据安全", 12, MUTED)
    footer(s, 7)


def s08_retrieval_chain(prs):
    s = blank(prs)
    bg(s)
    title(s, "神经扩散检索链路 ★", "技术核心 · 工程步骤与生物隐喻同一回事")
    flow = (
        "用户 Query\n"
        "  ↓ Embedding (BGE-M3)\n"
        "ANN 召回 Top-K 种子节点\n"
        "  ↓ Meta-path 白名单约束\n"
        "PPR 语义扩散 (mentions/contains/references/appears_in)\n"
        "  ↓ 子图聚合\n"
        "package → Skills 决策/生成/推荐"
    )
    txt(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.2), flow, 12, ACCENT2)
    table(s, Inches(6.6), Inches(1.55), Inches(6.0), Inches(4.8),
          ["工程", "生物隐喻", "一句话"],
          [
              ["Embedding", "感知", "信号编码"],
              ["ANN 种子", "动作电位", "少数神经元激活"],
              ["PPR 扩散", "神经传导", "沿突触传递"],
              ["边类型", "突触", "传导规则"],
              ["Meta-path", "皮层区", "语义分区"],
              ["Skills", "运动输出", "可执行行为"],
          ],
          [Inches(1.6), Inches(1.5), Inches(2.9)])
    txt(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.55),
        "算法演进：异构 GNN 研究 → 今日 PPR(静态 Message Passing) → L3 R-GCN/HAN(可学习)\n"
        "ANN 感知 · PPR 传导 · Meta-path 分区 · 边类型=突触规则", 10, GREEN)
    footer(s, 8)


def s09_rag_compare(prs):
    s = blank(prs)
    bg(s)
    title(s, "纯向量 RAG vs 谛听 ANN+PPR")
    table(s, Inches(0.6), Inches(1.6), Inches(12), Inches(2.2),
          ["", "纯向量 RAG", "谛听 ANN + PPR"],
          [
              ["检索", "只看语义相似", "相似 + 图结构多跳"],
              ["跨文档", "弱", "强 (appears_in / mentions)"],
              ["可解释", "chunk 列表", "路径 + 边 + node_id"],
              ["业务约束", "难", "Meta-path 白名单"],
          ],
          [Inches(2.2), Inches(4.9), Inches(4.9)])
    bullets(s, Inches(0.6), Inches(4.2), Inches(12), Inches(2.5), [
        "纯 LLM：闭卷考试，易幻觉",
        "纯向量 RAG：看不见结构关联（如旧报告失败案例 → 当前异常）",
        "ANN + PPR + Meta-path：向量负责「像不像」，图负责「连没连、该不该连」",
    ], 14)
    footer(s, 9)


def s10_arch(prs):
    s = blank(prs)
    bg(s)
    title(s, "系统架构：三层引擎")
    layers = [
        ("应用层", "谛听客户端：需求雷达 · 知识库 · 图谱 Tab · 认知推演 · 智能体对话\n服务端：Analyze · Ingestion · Graph API · Skill 编排"),
        ("本体层 ★", "节点=业务对象 · 边=语义关系 · 技能=可触发动作\n所有应用共享同一真相来源"),
        ("数据层", "Pipeline：分段入库 · HTML/图片 · OCR/VLM · 商城同步 · 上传 · 主动学习"),
    ]
    for i, (h, b) in enumerate(layers):
        y = Inches(1.6 + i * 1.65)
        card(s, Inches(1.2), y, Inches(10.8), Inches(1.4), h, b)
        if i < 2:
            txt(s, Inches(6.3), y + Inches(1.42), Inches(0.5), Inches(0.25), "▼", 14, ACCENT, align=PP_ALIGN.CENTER)
    footer(s, 10)


def s11_ontology(prs):
    s = blank(prs)
    bg(s)
    title(s, "知识网本体：系统的灵魂")
    txt(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.35), "传统数据模型 ✕", 14, ORANGE, True)
    bullets(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.5), [
        "✕ 数据在表行列中",
        "✕ 关系靠 JOIN 临时拼",
        "✕ 语义散落在 SQL/代码",
        "✕ 各应用看到不同「真相」",
    ], 12, MUTED)
    txt(s, Inches(6.5), Inches(1.6), Inches(6), Inches(0.35), "谛听知识网 ✓", 14, GREEN, True)
    bullets(s, Inches(6.5), Inches(2.0), Inches(6), Inches(2.5), [
        "✓ document → section → concept → media",
        "✓ 边持久化、有权重、有类型",
        "✓ 报告 → 段落 → 概念 → 谱图/商品图",
        "✓ 所有智能体读同一张网",
    ], 12, ACCENT2)
    txt(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.8),
        "分水岭：没有本体 LLM=通用助手；有了本体 LLM=领域专家", 15, WHITE, True)
    footer(s, 11)


def s12_flywheel(prs):
    s = blank(prs)
    bg(s)
    title(s, "知识网飞轮效应")
    steps = [("01 数据融合", "多源清洗抽取标准化"), ("02 语义成网", "实体挂节点 · meta-path 注册"),
             ("03 智能体赋能", "扩散子图驱动 Skill"), ("04 反馈强化", "边权更新 · 成功强化失败弱化")]
    for i, (h, b) in enumerate(steps):
        x = Inches(0.6 + i * 3.05)
        txt(s, x, Inches(2.0), Inches(0.5), Inches(0.4), h[:2], 22, ACCENT, True)
        txt(s, x + Inches(0.45), Inches(2.0), Inches(2.5), Inches(0.4), h[3:], 13, WHITE, True)
        txt(s, x, Inches(2.5), Inches(2.8), Inches(1.0), b, 11, LIGHT)
        if i < 3:
            txt(s, x + Inches(2.75), Inches(2.05), Inches(0.3), Inches(0.3), "→", 16, ACCENT)
    txt(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
        "用得越多 → 网越密 → 推荐越准 → 越离不开", 16, ACCENT2, True)
    footer(s, 12)


def s13_components(prs):
    s = blank(prs)
    bg(s)
    title(s, "三大构件：名词 · 关系 · 动词")
    comps = [
        ("节点 Nodes = 名词", "文档/段落/概念/媒体/商品/评价\n属性 + 向量 + 类型"),
        ("边 Links = 关系", "contains · references · mentions · appears_in\n持久 · 可遍历 · 第一公民"),
        ("技能 Skills = 动词", "检索/推荐/生成/挂接\n校验 · 来源 · 权限 · Human-in-the-Loop"),
    ]
    for i, (h, b) in enumerate(comps):
        x = Inches(0.6 + i * 4.1)
        card(s, x, Inches(1.7), Inches(3.9), Inches(2.5), h, b)
    txt(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.8),
        "mentions 关键：Section --mentions--> Concept\n连接非结构化载体与结构化概念 · PPR 跨文档传导主通道 · 生物别名「语义树突簇」", 11, ACCENT2)
    footer(s, 13)


def s14_ai_need(prs):
    s = blank(prs)
    bg(s)
    title(s, "为什么 AI 必须有知识网")
    bullets(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.5), [
        "没有知识网 ✕",
        "· 只能写通识摘要",
        "· 不知商品与概念关系",
        "· 无权控/无来源/无审计",
        "· 触发不了业务动作",
        "· 幻觉无法约束",
    ], 12, MUTED)
    bullets(s, Inches(6.6), Inches(1.6), Inches(6), Inches(3.5), [
        "有知识网 ✓",
        "· 理解完整业务语义",
        "· meta-path 检索真实数据",
        "· License/类目/角色权限",
        "· node_id + 来源可行动",
        "· 本体约束降幻觉",
    ], 12, ACCENT2)
    txt(s, Inches(0.6), Inches(5.3), Inches(12), Inches(1.2),
        "activate → seeds[]  ·  propagate → subgraph  ·  package → context  ·  llm → skill_output\n"
        "Skill 层不感知底层是 PPR 还是 GNN — 架构可演进", 11, LIGHT)
    footer(s, 14)


def s15_moat(prs):
    s = blank(prs)
    bg(s)
    title(s, "不可复制的护城河")
    moats = [
        ("01 数据网络效应", "每增数据源，关系网指数级丰富"),
        ("02 零边际成本", "本体已定，新场景=新视角"),
        ("03 组织知识固化", "规则/经验/失败案例编码进边权"),
        ("04 AI 倍增器", "LLM+专属知识网=领域专家"),
    ]
    for i, (h, b) in enumerate(moats):
        y = Inches(1.7 + i * 1.15)
        txt(s, Inches(0.6), y, Inches(0.5), Inches(0.4), h[:2], 18, ACCENT, True)
        txt(s, Inches(1.2), y, Inches(3), Inches(0.4), h[3:], 14, WHITE, True)
        txt(s, Inches(1.2), y + Inches(0.38), Inches(11), Inches(0.5), b, 12, LIGHT)
    footer(s, 15)


def s16_verify(prs):
    s = blank(prs)
    bg(s)
    title(s, "工程验证：从图谱到智能体")
    items = [
        ("L2 知识网基本盘 ✓", "beauty：数百商品 · 上千段落 · mentions/appears_in · PPR+meta_paths · verify_m21"),
        ("认知推演 ✓", "激活思考 → 子图召回 → 传导日志 → 思考回放 · 图谱 Tab 可遍历"),
        ("智能客服 POC ✓", "L2 子图 + L1 渐进对话 + L3 话术 Skill · 金牌导购非图书馆管理员"),
        ("需求雷达 ✓", "微信私域 → 类目过滤 → Chroma+LLM → 结构化需求 · 只读零操作"),
    ]
    for i, (h, b) in enumerate(items):
        y = Inches(1.65 + i * 1.25)
        txt(s, Inches(0.6), y, Inches(3.5), Inches(0.35), h, 13, ACCENT, True)
        txt(s, Inches(0.6), y + Inches(0.35), Inches(12), Inches(0.7), b, 11, LIGHT)
    txt(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.35),
        "机制已通 → 下一页：技能集四类智能体应用（客服/写作/短视频/直播）", 12, GREEN)
    footer(s, 16)


def s17_skill_apps(prs):
    s = blank(prs)
    bg(s)
    title(s, "技能集 · 智能体应用案例", "商业垂直 · AI 智能 · 技能集 · 同网不同 Skill")
    cols = [
        ("商业垂直", "被动投喂 + 主动学习", "beauty 960 条 · 防晒/护肤/彩妆…"),
        ("AI 智能", "引擎与技术栈", "文本/图像/视频 · 大模型连接"),
        ("技能集", "策略包 + 版本", "四类 Skill · 人工维护"),
    ]
    for i, (h, sub, b) in enumerate(cols):
        card(s, Inches(0.6 + i * 4.1), Inches(1.55), Inches(3.9), Inches(1.25), h, f"{sub}\n{b}", ACCENT if i == 2 else ACCENT2)
    table(s, Inches(0.6), Inches(3.0), Inches(12), Inches(2.55),
          ["技能集", "应用案例", "策略", "状态"],
          [
              ["智能客服", "电商售前导购", "cs_dialogue · PDS 四步", "✅ M6-1/2"],
              ["智能写作", "文案写手", "doc_writer · 信任状脚手架", "✅ M6-3"],
              ["智能视频制作", "短视频创作", "编剧/评估/美术 Skill", "🚧 M7-1"],
              ["智能数字人直播", "数字人直播", "live_stream（预留）", "🔜 远期"],
          ],
          [Inches(2.2), Inches(2.4), Inches(4.0), Inches(1.8)])
    txt(s, Inches(0.6), Inches(5.75), Inches(12), Inches(0.35),
        "activate → propagate → package → llm → skill_output  ·  四类产品共用契约", 11, GREEN)
    txt(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.35),
        "不是各建一套 AI — 在一张网上用不同 Skill 长出四只「手」", 12, ACCENT2, True)
    footer(s, 17)


def s18_scenarios(prs):
    s = blank(prs)
    bg(s)
    title(s, "垂直扩展", "换 category 不换架构")
    sc = [
        ("私域运营", "需求雷达 · 结构化需求沉淀 · 主链路"),
        ("研发知识", "工艺/专利/谱图 · meta-path 经验传导"),
        ("企业培训", "认知推演 · 附来源的经验传导"),
        ("生鲜/家电", "cs_fresh / cs_appliance 话术变体"),
        ("更多垂直", "本体一次建设 · Skill 按场景叠加"),
    ]
    for i, (h, b) in enumerate(sc):
        col, row = i % 3, i // 3
        card(s, Inches(0.6 + col * 4.1), Inches(1.7 + row * 2.3), Inches(3.9), Inches(2.0), h, b)
    txt(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.35),
        "beauty 为样板 · 机制一次建设 · 垂直按类目扩展", 13, ACCENT2, True)
    footer(s, 18)


def s18_rag(prs):
    s = blank(prs)
    bg(s)
    title(s, "RAG 与模型选型", "检索增强生成 · 开卷考试而非闭卷瞎编")
    table(s, Inches(0.6), Inches(1.55), Inches(12), Inches(2.0),
          ["环节", "常规定义", "谛听实现"],
          [
              ["检索端", "向量 Top-K chunk", "ANN + PPR + Meta-path → 结构化子图"],
              ["生成端", "LLM 读 chunk", "读知识网片段 → 话术/建议/报告"],
              ["价值", "降幻觉", "私有沉淀 · 来源边权可查"],
          ],
          [Inches(1.5), Inches(4.0), Inches(6.5)])
    bullets(s, Inches(0.6), Inches(3.8), Inches(12), Inches(2.8), [
        "LLM：DeepSeek 等 · 必须被知识网约束",
        "向量：BGE-M3 + Chroma · ANN 种子",
        "图扩散：PPR + meta_paths 白名单",
        "VLM：media 节点描述 · 跨文档反查",
        "小模型：XGBoost/LightGBM + 符号规则混合推理",
        "原则：RAG(ANN+PPR)优先 · 小模型专精 · Skill 调度 · 不盲目微调",
    ], 11, LIGHT)
    footer(s, 19)


def s19_paradigm(prs):
    s = blank(prs)
    bg(s)
    title(s, "范式跃迁：三代进化")
    eras = [
        ("1.0 肌肉型", "工具各自为政 · 响应以月计"),
        ("2.0 感知型", "BI 看板 · 看到≠行动 · 响应以周计"),
        ("3.0 神经系统", "感知→理解→思考→行动→反馈 · 响应以小时/天计"),
    ]
    for i, (h, b) in enumerate(eras):
        x = Inches(0.6 + i * 4.1)
        card(s, x, Inches(1.7), Inches(3.9), Inches(2.5), h, b)
        if i < 2:
            txt(s, x + Inches(3.85), Inches(2.3), Inches(0.3), Inches(0.4), "→", 18, ACCENT)
    txt(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
        "数字孪生不再是静态地图，而是活的控制系统", 14, ACCENT2, True)
    footer(s, 20)


def s20_nervous_rise(prs):
    s = blank(prs)
    bg(s)
    title(s, "企业正在长出「神经系统」")
    bullets(s, Inches(0.6), Inches(1.6), Inches(12), Inches(4.5), [
        "按席位 SaaS → 按价值消耗的智能体（买 Decision 非 Feature）",
        "功能驱动路线图 → 本体驱动（先语义，再应用）",
        "人工中间层信息传递 → 知识网+AI 压缩",
        "静态 BI 报表 → 系统代我行动（保留 Human-in-the-Loop）",
        "",
        "终局：竞争不只是数据多，而是谁的神经系统更完整",
    ], 14)
    footer(s, 21)


def s21_loop(prs):
    s = blank(prs)
    bg(s)
    title(s, "干湿闭环：从地图到自动驾驶仪")
    steps = [
        ("1 感知", "Pipeline 多模态入库"),
        ("2 理解", "实体抽取 · 语义节点挂接"),
        ("3 推理", "Agent 扩散 · 缺口检测 · 场景模拟"),
        ("4 行动", "话术/方案/挂链 · 人审后执行"),
        ("5 学习", "反馈回流 · 边权更新 · 突触可塑性 L3"),
    ]
    for i, (h, b) in enumerate(steps):
        y = Inches(1.6 + i * 0.95)
        txt(s, Inches(0.6), y, Inches(1.3), Inches(0.35), h, 13, ACCENT, True)
        txt(s, Inches(2.0), y, Inches(10.5), Inches(0.7), b, 12, LIGHT)
    txt(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
        "生物侧：可塑性/传导/反馈  ↔  数字侧：图/向量/Agent/Pipeline", 12, ACCENT2)
    footer(s, 22)


def s22_plasticity(prs):
    s = blank(prs)
    bg(s)
    title(s, "突触可塑性：会学习的知识网", "L3 可塑层")
    items = [
        ("强化", "检索命中/Analyze 采纳 → 边权上升"),
        ("退化", "长期未命中/低置信 → 软隐藏/待清理"),
        ("Hub 增厚", "appears_in 枢纽 · PPR 停留概率提高"),
        ("工程字段", "props_json: weight · traversal_count · staleness · pruned"),
    ]
    for i, (h, b) in enumerate(items):
        y = Inches(1.7 + i * 1.15)
        txt(s, Inches(0.6), y, Inches(1.5), Inches(0.35), h, 14, ACCENT, True)
        txt(s, Inches(2.2), y, Inches(10), Inches(0.7), b, 12, LIGHT)
    footer(s, 23)


def s23_decision(prs):
    s = blank(prs)
    bg(s)
    title(s, "终局：Decision Intelligence")
    bullets(s, Inches(0.6), Inches(1.6), Inches(12), Inches(4.5), [
        "组织变网状：一线与 Agent 协作，管理层聚焦战略",
        "跨组织知识互联：产业链级本体互认",
        "软件形态变：建操作系统，非买更多模块",
        "决策速度 = 降维打击：分钟级 vs 周级",
        "",
        "问题不是「要不要上 AI」，而是「准备好拥有自己的知识网了吗？」",
    ], 14)
    footer(s, 24)


def s24_three(prs):
    s = blank(prs)
    bg(s)
    title(s, "三句话带走")
    items = [
        ("01", "构建知识网本体，而非堆叠 Feature", "统一语义层 · 护城河非功能清单"),
        ("02", "让 AI 连接业务流程，而非停在 Chat", "价值在行动化，非对话化"),
        ("03", "用操作系统思维取代工具思维", "一平台无限 Skill · 边际成本趋零"),
    ]
    for i, (n, h, b) in enumerate(items):
        y = Inches(1.7 + i * 1.55)
        txt(s, Inches(0.6), y, Inches(0.6), Inches(0.5), n, 22, ACCENT, True)
        txt(s, Inches(1.3), y, Inches(11), Inches(0.4), h, 15, WHITE, True)
        txt(s, Inches(1.3), y + Inches(0.42), Inches(11), Inches(0.7), b, 12, LIGHT)
    footer(s, 25)


def s25_closing(prs):
    s = blank(prs)
    bg(s)
    txt(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
        "未来属于那些「看得见」的人", 40, WHITE, True, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.8),
        "看得见关系 · 看得见传导路径 · 看得见 AI 该行动在哪里", 16, ACCENT2, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
        "基于生物-数字双模态知识网与神经扩散机制的智能体系统研发及应用", 14, LIGHT, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.45),
        "客服 · 写作 · 短视频 · 数字人直播 — 同一张网上的四只「手」", 12, MUTED, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
        "谛听 DITING · 感谢聆听", 14, ACCENT, align=PP_ALIGN.CENTER)
    footer(s, 26)


# ── Appendix slides (速查，不讲) ─────────────────────────

def appendix_divider(prs, letter, title_text):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    txt(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5), f"附录 {letter}", 14, MUTED)
    txt(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.8), title_text, 36, ACCENT, True)
    txt(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5), "速查资料 · 宣讲时不讲", 14, ORANGE, align=PP_ALIGN.CENTER)
    return s


def a_qa(prs, num):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, "附录 A · 高频 Q&A", "速查 · 不讲")
    qa1 = [
        ("Q：和 Palantir / 传统知识图谱有何不同？",
         "A：Palantir 卖 Ontology OS；我们用生物隐喻定义异构图，PPR 已落地，私域场景验证闭环。更轻、可私有化、Skill 可插拔。"),
        ("Q：生物隐喻是噱头吗？",
         "A：每层隐喻映射工程对象；早期起点是异构 GNN+知识图谱 Message Passing，PPR 是其可部署版本。"),
        ("Q：ANN 和 PPR 为何不能只用其一？",
         "A：ANN 秒级找起点（感知）；PPR 沿结构深化（传导）。只做 ANN 漏结构关键节点；只做 PPR 无种子。"),
    ]
    qa2 = [
        ("Q：Meta-path 是什么？",
         "A：异构图类型化路径模板。按任务选 MP1/MP2/MP5，避免乱游走。"),
        ("Q：一定要微调大模型吗？",
         "A：不一定。默认 RAG+知识网约束。本体质量 > 模型规模。"),
        ("Q：GNN 做了吗？和 PPR 什么关系？",
         "A：研究渊源在异构 GNN Message Passing；生产默认 PPR；L3 可上 R-GCN/HAN，接口不变。"),
        ("Q：部署与数据安全？",
         "A：端云分离 · 私有化 · category 隔离 · License · 敏感数据不出域。"),
    ]
    items = qa1 if num == 1 else qa2
    step = 1.45 if num == 2 else 1.75
    for i, (q, a) in enumerate(items):
        y = Inches(1.5 + i * step)
        txt(s, Inches(0.6), y, Inches(12), Inches(0.45), q, 12, ACCENT, True)
        txt(s, Inches(0.6), y + Inches(0.42), Inches(12), Inches(0.9), a, 11, LIGHT)
    footer(s, APPENDIX_START + num - 1, "附录速查")


def b_overview(prs):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, "附录 B · 术语总览", "专业定义 + 易懂版")
    table(s, Inches(0.6), Inches(1.5), Inches(12), Inches(4.8),
          ["术语", "专业一句话", "易懂版"],
          [
              ["RAG", "先检索外部知识再 LLM 生成", "开卷考试"],
              ["ANN", "高维近似最近邻检索", "感官定位"],
              ["PPR", "种子出发的个性化 PageRank", "神经传导"],
              ["Meta-path", "异构图类型化路径模板", "皮层功能区"],
              ["mentions", "文本/媒体提及概念/实体", "语义树突"],
              ["神经扩散", "ANN+PPR 沿 meta-path 传播", "感知+传导"],
              ["异构 GNN", "HIN 上 Message Passing 研究起点", "算法渊源"],
          ],
          [Inches(1.8), Inches(5.2), Inches(5.0)])
    footer(s, APPENDIX_START + 2, "附录速查")


def b_term(prs, num, heading, blocks):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, f"附录 B · {heading}", "术语速查")
    y = Inches(1.5)
    for label, content in blocks:
        txt(s, Inches(0.6), y, Inches(12), Inches(0.35), label, 12, ACCENT, True)
        txt(s, Inches(0.6), y + Inches(0.32), Inches(12), Inches(1.5), content, 10, LIGHT)
        y += Inches(0.32 + min(1.6, 0.25 + len(content) / 80))
    footer(s, APPENDIX_START + 2 + num, "附录速查")


def b_algorithm_evolution(prs):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, "附录 B · 算法演进", "异构 GNN 研究 → PPR 工程 → L3 可学习扩散")
    table(s, Inches(0.6), Inches(1.5), Inches(12), Inches(2.0),
          ["阶段", "算法栈", "状态"],
          [
              ["研究起点", "异构图神经网络(HIN/GNN) + KG Message Passing", "算法渊源"],
              ["L2 传导", "ANN + PPR + meta-path（静态 Message Passing）", "✅ 生产默认"],
              ["L3 可塑", "R-GCN / HAN（可学习 Message Passing）", "🔜 远期可选"],
          ],
          [Inches(1.8), Inches(7.2), Inches(3.0)])
    bullets(s, Inches(0.6), Inches(3.8), Inches(12), Inches(2.8), [
        "为何先 PPR：2GB 无 GPU · 冷启动可解释 · 路径/边/meta-path 可审计",
        "PPR = 无需训练参数的 Message Passing（入库方案 §14.3.2）",
        "统一接口：activate → propagate → package → llm",
        "  今日 ppr_expand() · 明日 gnn_forward() · Skill 层不感知",
        "宣讲口径：有算法纵深 · 今日交付 PPR · 明日预留 GNN",
    ], 11, LIGHT)
    footer(s, APPENDIX_START + 8, "附录速查")


def c_index(prs):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, "附录 C · PPT 页码对照")
    rows = [
        ["1", "封面"], ["2", "数据黑盒"], ["3", "传统困局"], ["4", "我们的方向"],
        ["5", "核心理念"], ["6", "谛听 vs 传统"], ["7", "三大能力层"], ["8", "★神经扩散链路"],
        ["9", "RAG 对比"], ["10", "系统架构"], ["11", "知识网本体"], ["12", "飞轮"],
        ["13", "节点/边/技能"], ["14", "AI 需知识网"], ["15", "护城河"], ["16", "工程验证"],
        ["17", "★技能集应用"], ["18", "垂直扩展"], ["19", "RAG 选型"], ["20", "范式跃迁"],
        ["21", "神经系统"], ["22", "干湿闭环"], ["23", "突触可塑性"], ["24", "Decision Intel"],
        ["25", "三句话"], ["26", "结语"],
    ]
    half = len(rows) // 2 + 1
    for i, (p, t) in enumerate(rows[:half]):
        txt(s, Inches(0.6), Inches(1.5 + i * 0.22), Inches(5.5), Inches(0.22), f"{p:>2}  {t}", 10, LIGHT)
    for i, (p, t) in enumerate(rows[half:]):
        txt(s, Inches(6.5), Inches(1.5 + i * 0.22), Inches(5.5), Inches(0.22), f"{p:>2}  {t}", 10, LIGHT)
    footer(s, APPENDIX_START + 9, "附录速查")


def d_cheatsheet(prs):
    s = blank(prs)
    bg(s, APPENDIX_BG)
    title(s, "附录 D · 一页纸备忘", "上台前 30 秒复习")
    txt(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
        "核心命题：基于生物-数字双模态知识网与神经扩散机制的智能体系统研发及应用", 12, ACCENT2, True)
    txt(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.3),
        "（文档短标题：生物-数字双模态知识网机制的智能体研发与应用）", 10, MUTED)
    txt(s, Inches(0.6), Inches(2.25), Inches(12), Inches(0.35),
        "四类应用：电商售前导购 ✅ · 文案写手 ✅ · 短视频 🚧 · 数字人直播 🔜", 11, GREEN)
    flow = (
        "Query → Embedding → ANN 种子 → Meta-path + PPR → 子图 → Skills\n"
        "感知     编码        动作电位      皮层区+神经传导    读出    行动"
    )
    txt(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.0), flow, 12, LIGHT)
    bullets(s, Inches(0.6), Inches(3.5), Inches(12), Inches(2.5), [
        "RAG = 开卷考试（检索端 = ANN+PPR，非普通 chunk）",
        "PPR = 从种子沿边扩散打分（神经传导）",
        "Meta-path = 按业务选路径模板（皮层功能区）",
        "异构 GNN = 算法渊源 · 今日 PPR · L3 可上 R-GCN/HAN",
        "",
        "一句差异化：向量找「像不像」，图扩散找「连没连、该不该连」",
    ], 13, WHITE)
    footer(s, APPENDIX_START + 10, "附录速查")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    main = [
        s01_cover, s02_dark_age, s03_dilemma, s04_direction, s05_philosophy,
        s06_compare, s07_layers, s08_retrieval_chain, s09_rag_compare, s10_arch,
        s11_ontology, s12_flywheel, s13_components, s14_ai_need, s15_moat,
        s16_verify, s17_skill_apps, s18_scenarios, s18_rag, s19_paradigm, s20_nervous_rise,
        s21_loop, s22_plasticity, s23_decision, s24_three, s25_closing,
    ]
    for fn in main:
        fn(prs)

    appendix_divider(prs, "A–D", "术语速查 · Q&A · 页码 · 备忘")

    a_qa(prs, 1)
    a_qa(prs, 2)
    b_overview(prs)

    b_term(prs, 1, "RAG", [
        ("定义", "检索增强生成：用户提问时先从外部知识库检索，再将检索结果+原问题交给 LLM。"),
        ("谛听", "检索端=ANN+PPR+Meta-path→结构化子图；生成端=读知识网片段；价值=开卷考试、来源可查。"),
    ])
    b_term(prs, 2, "PPR", [
        ("定义", "从种子节点出发的个性化 PageRank 随机游走，计算局部关联强度。"),
        ("公式", "π = α·P·π + (1−α)·v  （v=种子个性化向量，α=扩散阻尼）"),
        ("谛听", "种子来自 ANN；沿 mentions/contains 多跳扩散；模拟神经传导/激活扩散。"),
    ])
    b_term(prs, 3, "Meta-path", [
        ("定义", "异构图上 A₁─R₁→A₂─R₂→… 的类型化关系路径模板。"),
        ("示例", "MP1: concept→section←document | MP2: section→media→document | MP5: media→…→document"),
        ("谛听", "PPR 仅沿 meta_paths 白名单游走；按任务选 path，保证业务逻辑。"),
    ])
    b_term(prs, 4, "ANN + PPR", [
        ("ANN", "BGE 空间 Top-K · 秒级定位起点 · 生物映射=感知/感官编码"),
        ("PPR", "以种子图上游走 · 多跳传导 · 生物映射=神经传导/激活场"),
        ("联合", "ANN 广撒网 + PPR 深化 + Meta-path 分区 + 边类型=突触规则"),
    ])
    b_term(prs, 5, "Mentions & 边类型", [
        ("mentions", "文本/媒体提及概念/实体 · 语义树突簇 · 跨文档 PPR 主通道"),
        ("边速查", "contains=胞体胞突 | references=短程突触 | appears_in=长程髓鞘 | related_product=挂接轴突"),
    ])

    b_algorithm_evolution(prs)

    c_index(prs)
    d_cheatsheet(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)} (宣讲 {MAIN_COUNT} + 附录 {len(prs.slides) - MAIN_COUNT})")


if __name__ == "__main__":
    build()
