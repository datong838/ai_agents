#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成湃肽生物（ZPC）合作方案 PPTX。

输出：docs/ZPC-湃肽生物-多肽研发智能知识网络方案.pptx
方案：docs/M7-2-ZPC-湃肽生物合作方案PPT.md
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── 主题常量 ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG = RGBColor(0x0A, 0x16, 0x28)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xD6, 0xE0)
MUTED = RGBColor(0x8A, 0x9B, 0xAD)
FONT = "Microsoft YaHei"

OUTPUT = (
    Path(__file__).resolve().parent.parent
    / "ZPC-湃肽生物-多肽研发智能知识网络方案.pptx"
)


# ── 辅助函数 ──────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor = BG) -> None:
    """设置幻灯片纯色背景。"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 18,
    color: RGBColor = WHITE,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT,
) -> None:
    """在幻灯片上添加文本框。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def add_bullet_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    font_size: int = 16,
    color: RGBColor = LIGHT,
    line_spacing: float = 1.3,
) -> None:
    """添加多行要点文本框。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
        p.level = 0
    return box


def add_footer(slide, page_num: int, subtitle: str = "谛听 × 湃肽生物") -> None:
    """添加页脚与页码。"""
    add_textbox(
        slide,
        Inches(0.6),
        Inches(7.0),
        Inches(8),
        Inches(0.35),
        subtitle,
        font_size=10,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(12.2),
        Inches(7.0),
        Inches(0.8),
        Inches(0.35),
        str(page_num),
        font_size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_accent_line(slide, left, top, width) -> None:
    """添加青色装饰横线。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_slide_title(slide, title: str, page_num: int) -> None:
    """标准内页标题区。"""
    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(12),
        Inches(0.7),
        title,
        font_size=28,
        color=WHITE,
        bold=True,
    )
    add_accent_line(slide, Inches(0.6), Inches(1.15), Inches(2.2))
    add_footer(slide, page_num)


def _style_table_cell(cell, text: str, font_size: int = 12, bold: bool = False,
                      color: RGBColor = LIGHT, align=PP_ALIGN.LEFT) -> None:
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = BG
    for p in cell.text_frame.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, left, top, width, height, headers, rows,
              col_widths=None, header_color=ACCENT) -> None:
    """添加主题风格表格。"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        _style_table_cell(table.cell(0, j), h, font_size=13, bold=True,
                          color=header_color, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            bold = j == 0
            _style_table_cell(table.cell(i + 1, j), val, font_size=11, bold=bold)


# ── 24 张幻灯片 ───────────────────────────────────────────
def slide_01_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0.6), Inches(2.55), Inches(4.5))
    add_textbox(slide, Inches(0.6), Inches(1.6), Inches(10), Inches(0.9),
                "谛听  DITING", font_size=44, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.6), Inches(2.75), Inches(11), Inches(0.8),
                "让每一条多肽工艺都被「看见」", font_size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.6), Inches(3.65), Inches(11), Inches(0.6),
                "生物-数字双模态知识网 × 神经扩散智能体", font_size=22, color=LIGHT)
    add_textbox(slide, Inches(0.6), Inches(5.8), Inches(11), Inches(0.5),
                "湃肽生物（ZPC）研发中心合作方案", font_size=18, color=MUTED)
    add_textbox(slide, Inches(0.6), Inches(6.35), Inches(11), Inches(0.4),
                "2026 · 谛听科技", font_size=12, color=MUTED)


def slide_02_partner(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "合作方 · 湃肽生物（ZPC）", 2)
    add_bullet_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), [
        "国家级专精特新「小巨人」企业",
        "双轮驱动：药物肽 + 化妆品肽",
        "cGMP / FDA 双认证生产体系",
        "浙江省多肽生物工程研究中心",
        "核心技术：固相合成 · 液相合成 · 环肽 · 修饰肽",
        "在研管线：司美格鲁肽、利拉鲁肽、替尔泊肽等",
    ], font_size=15)
    add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5),
                "核心痛点", font_size=18, color=ACCENT, bold=True)
    add_bullet_textbox(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5), [
        "20 年工艺经验沉淀在专家脑中，难以传承",
        "LIMS / ELN / 文献 / 专利数据各自为政",
        "新人培养周期长，关键决策高度依赖资深工程师",
        "AI 工具缺乏领域本体约束，难以进入核心研发流程",
        "实验数据与知识资产未形成可复用的智能网络",
    ], font_size=14, color=LIGHT)


def slide_03_blackbox(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "多肽研发「黑盒时代」", 3)
    metrics = [
        ("20+", "年工艺经验沉淀"),
        ("70%", "关键决策依赖专家直觉"),
        ("5+", "异构数据源孤岛并存"),
    ]
    for i, (num, label) in enumerate(metrics):
        x = Inches(0.8 + i * 4.1)
        add_textbox(slide, x, Inches(2.0), Inches(3.5), Inches(1.0),
                    num, font_size=56, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.1), Inches(3.5), Inches(0.6),
                    label, font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)
    add_bullet_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(2.5), [
        "工艺卡片、实验记录、HPLC 谱图、专利文献——信息丰富但彼此割裂",
        "资深工程师退休或转岗，隐性知识随之流失，组织记忆出现断层",
        "通用大模型无法「理解」多肽领域的因果链：参数 → 收率 → 杂质谱",
        "研发效率的瓶颈不在算力，而在于知识不可见、不可传导、不可行动化",
    ], font_size=15)


def slide_04_silos(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "传统信息化困局：数据孤岛", 4)
    silos = [
        ("LIMS", "实验流程管理", "批次数据难关联工艺知识"),
        ("ELN", "电子实验记录", "记录与检索分离，难成网"),
        ("专利库", "知识产权资产", "全文难结构化，难迁移"),
        ("文献库", "外部知识输入", "与内部工艺缺乏映射"),
    ]
    for i, (name, sub, pain) in enumerate(silos):
        x = Inches(0.5 + (i % 2) * 6.4)
        y = Inches(1.6 + (i // 2) * 2.6)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y, Inches(5.9), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.5),
                    name, font_size=22, color=ACCENT, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.75), Inches(5.3), Inches(0.4),
                    sub, font_size=14, color=MUTED)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.2), Inches(5.3), Inches(0.8),
                    pain, font_size=14, color=LIGHT)
    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
                "→ 系统越多，孤岛越多；知识依然「看不见」", font_size=14, color=ACCENT)


def slide_05_positioning(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "我们的定位：L1 / L2 / L3 传导架构", 5)
    layers = [
        ("L1  数据管道层", "采集 · 解析 · 结构化入库",
         "工艺卡片 / 实验记录 / 谱图 / 专利 / 文献 → 统一语义切片"),
        ("L2  知识网本体层", "节点 · 边 · 突触权重",
         "document / section / concept / media / experiment 映射多肽世界"),
        ("L3  智能体行动层", "PPR 扩散 + Skill 调度",
         "纯化推荐 · 路线选型 · 失效诊断 · 人机协同决策"),
    ]
    for i, (title, sub, desc) in enumerate(layers):
        y = Inches(1.5 + i * 1.85)
        add_accent_line(slide, Inches(0.6), y, Inches(0.08))
        add_textbox(slide, Inches(0.85), y - Inches(0.05), Inches(4), Inches(0.45),
                    title, font_size=20, color=ACCENT, bold=True)
        add_textbox(slide, Inches(0.85), y + Inches(0.4), Inches(4), Inches(0.35),
                    sub, font_size=13, color=MUTED)
        add_textbox(slide, Inches(5.2), y, Inches(7.5), Inches(1.5),
                    desc, font_size=15, color=LIGHT)
    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.35),
                "生物-数字双模态知识网：让数据流动，让经验传导，让决策可行动", font_size=14, color=ACCENT)


def slide_06_philosophy(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "核心理念", 6)
    pillars = [
        ("增强专家，而非替代", "AI 是资深工程师的「第二大脑」，每条建议附可追溯来源"),
        ("知识网本体先行", "先建业务语义层，再谈大模型；本体是 AI 落地的地基"),
        ("人机协同闭环", "Human-in-the-Loop：专家审核 → 反馈 → 边权强化"),
        ("数据 → 决策 → 实验", "从被动检索升级为主动推荐，驱动下一轮实验设计"),
    ]
    for i, (title, desc) in enumerate(pillars):
        y = Inches(1.55 + i * 1.35)
        add_textbox(slide, Inches(0.6), y, Inches(4.5), Inches(0.45),
                    f"0{i+1}  {title}", font_size=18, color=ACCENT, bold=True)
        add_textbox(slide, Inches(0.6), y + Inches(0.45), Inches(12), Inches(0.7),
                    desc, font_size=15, color=LIGHT)


def slide_07_comparison(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "谛听方案 vs 传统信息化", 7)
    headers = ["维度", "传统方案", "谛听知识网方案"]
    rows = [
        ["数据模型", "表格 + 全文检索", "节点/边/技能 本体图谱"],
        ["知识传承", "文档归档，依赖人工翻阅", "PPR 扩散，相似工艺自动迁移"],
        ["AI 接入", "通用 ChatBot，易幻觉", "meta-path 约束 + 来源溯源"],
        ["专家角色", "AI 替代焦虑", "增强专家，审核即训练"],
        ["实验闭环", "记录与分析割裂", "干湿闭环，实验回流强化边权"],
        ["部署模式", "SaaS 数据出境", "私有化部署，数据不出域"],
    ]
    add_table(slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.2),
              headers, rows, col_widths=[Inches(2.2), Inches(4.8), Inches(5.1)])


def slide_08_capabilities(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "三大能力层", 8)
    caps = [
        ("知识网", "Knowledge Network",
         "多肽领域专属本体：文档、段落、概念、谱图、实验、分子节点互联"),
        ("传导引擎", "PPR Diffusion Engine",
         "个性化 PageRank 沿 meta-path 扩散，模拟神经突触信号传导"),
        ("领域智能体", "Domain Agents",
         "纯化选型 · 合成路线 · 失效诊断 · 文献检索 · 方案生成"),
    ]
    for i, (cn, en, desc) in enumerate(caps):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.7), Inches(3.9), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, x + Inches(0.25), Inches(1.9), Inches(3.4), Inches(0.5),
                    cn, font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.25), Inches(2.45), Inches(3.4), Inches(0.4),
                    en, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.25), Inches(3.1), Inches(3.4), Inches(3.0),
                    desc, font_size=14, color=LIGHT)


def slide_09_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "三层架构：应用 / 本体 / 管道", 9)
    layers = [
        ("应用层", "对话检索 · 智能推荐 · 实验方案生成 · 知识图谱可视化",
         "研发人员日常交互入口，专家审核与反馈"),
        ("本体层", "节点 / 边 / 技能 / meta-path / 突触权重",
         "多肽业务语义核心，连接数据与行动的桥梁"),
        ("管道层", "采集 → 解析 → 分段 → 向量化 → 成网 → 回流",
         "LIMS/ELN/文献/专利/谱图统一入库与持续更新"),
    ]
    for i, (name, items, note) in enumerate(layers):
        y = Inches(1.55 + i * 1.85)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.6), y, Inches(12.1), Inches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, Inches(0.85), y + Inches(0.15), Inches(2.5), Inches(0.45),
                    name, font_size=20, color=ACCENT, bold=True)
        add_textbox(slide, Inches(3.5), y + Inches(0.15), Inches(8.8), Inches(0.45),
                    items, font_size=14, color=WHITE)
        add_textbox(slide, Inches(3.5), y + Inches(0.7), Inches(8.8), Inches(0.6),
                    note, font_size=12, color=MUTED)


def slide_10_ontology_vs_model(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "知识网本体 vs 传统数据模型", 10)
    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.45),
                "传统数据模型", font_size=18, color=MUTED, bold=True)
    add_bullet_textbox(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), [
        "关系型表：行与列，缺乏语义",
        "全文索引：关键词匹配，无语境",
        "文件夹：树状层级，跨文档断裂",
        "向量库：语义相近但因果不明",
        "→ 知道「有什么」，不知道「为什么」",
    ], font_size=14, color=MUTED)
    add_textbox(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(0.45),
                "知识网本体", font_size=18, color=ACCENT, bold=True)
    add_bullet_textbox(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), [
        "节点：document / section / concept / media / experiment",
        "边：contains / references / similar_to / leads_to / failed_with",
        "技能：纯化推荐、路线选型、失效诊断等可触发操作",
        "meta-path：约束 PPR 扩散路径，保障领域逻辑",
        "→ 理解结构、因果与行动，支撑 AI 决策",
    ], font_size=14, color=LIGHT)


def slide_11_flywheel(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "飞轮效应：知识网持续强化", 11)
    steps = [
        ("01  入库", "工艺卡片 / 实验 / 谱图结构化入库"),
        ("02  成网", "节点互联，边权初始化，图谱可浏览"),
        ("03  智能体", "对话检索 + 推荐 + 诊断，专家日常使用"),
        ("04  反馈强化", "审核采纳 → 边权更新 → 网络更精准"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 3.15)
        # arrow circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.0),
                                        Inches(1.4), Inches(1.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        circle.line.color.rgb = ACCENT
        add_textbox(slide, x + Inches(0.9), Inches(2.35), Inches(1.4), Inches(0.7),
                    f"{i+1}", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.6), Inches(3.1), Inches(0.45),
                    title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(4.1), Inches(3.1), Inches(1.2),
                    desc, font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(slide, x + Inches(2.5), Inches(2.5), Inches(0.5), Inches(0.4),
                        "→", font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
                "每多一条实验回流、每一次专家采纳，知识网就变得更聪明——这是不可速成的护城河",
                font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_12_three_components(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "三构件：节点 / 边 / 技能", 12)
    components = [
        ("节点 Nodes", "神经元",
         ["document — 工艺卡片、专利、文献",
          "section — 纯化步骤、合成路线段",
          "concept — 序列、树脂、保护基、收率",
          "media — HPLC、质谱、电镜谱图",
          "experiment — 批次实验三元组"]),
        ("边 Links", "突触",
         ["contains — 文档→段落→谱图",
          "references — 段落引用色谱图",
          "similar_to — 相似序列/工艺",
          "leads_to — 参数→收率因果",
          "failed_with — 失败路径抑制"]),
        ("技能 Skills", "动作电位",
         ["纯化参数推荐",
          "合成路线选型",
          "失效诊断排查",
          "文献/工艺检索",
          "实验方案生成"]),
    ]
    for i, (title, alias, items) in enumerate(components):
        x = Inches(0.5 + i * 4.2)
        add_textbox(slide, x, Inches(1.55), Inches(3.9), Inches(0.45),
                    title, font_size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.0), Inches(3.9), Inches(0.35),
                    alias, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_bullet_textbox(slide, x, Inches(2.5), Inches(3.9), Inches(4.2),
                           [f"• {it}" for it in items], font_size=12, color=LIGHT)


def slide_13_why_ontology(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "为什么 AI 需要知识网？", 13)
    add_bullet_textbox(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(5.0), [
        "通用 RAG 只有「语义相近」，没有「业务因果」——无法回答「为什么这次收率低」",
        "多肽领域知识高度结构化：序列 → 保护基 → 偶联 → 纯化 → 收率/杂质，需要图谱承载",
        "meta-path 约束 PPR 扩散：只在「合法路径」上传播，杜绝跨领域幻觉",
        "每个推荐可溯源：答案附带工艺卡片段落、谱图、实验批次，专家可一键核验",
        "技能（Skill）将知识转化为行动：不只「告诉」，更能「推荐参数、生成方案」",
        "知识网随使用生长：专家反馈即训练信号，无需昂贵全量微调",
    ], font_size=16)


def slide_14_moat(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "护城河：不可复制的四重壁垒", 14)
    moats = [
        ("湃肽专属工艺图谱", "20 年积累的结构化映射，通用模型无法复制"),
        ("持续可塑的突触权重", "每次实验回流与专家采纳都在强化网络"),
        ("领域 meta-path 模板", "多肽因果链路径是领域 Know-how 的结晶"),
        ("人机协同信任链", "来源溯源 + 专家审核，建立组织级 AI 信任"),
    ]
    for i, (title, desc) in enumerate(moats):
        y = Inches(1.55 + i * 1.35)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), y, Inches(12.1), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, Inches(0.9), y + Inches(0.15), Inches(4.5), Inches(0.4),
                    f"▸  {title}", font_size=17, color=ACCENT, bold=True)
        add_textbox(slide, Inches(5.5), y + Inches(0.2), Inches(7.0), Inches(0.7),
                    desc, font_size=14, color=LIGHT)


def slide_15_scenarios(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "三大核心场景", 15)
    scenarios = [
        ("纯化选型", "输入氨基酸序列",
         "推荐柱型 / 填料 / 梯度 / 流速，附相似成功案例与谱图"),
        ("合成路线", "长肽 / 环肽 / 修饰肽",
         "分段合成 vs 全长合成，保护基策略，偶联方案对比"),
        ("失效诊断", "收率低 / 新杂质 / 聚集",
         "沿 failed_with 边排查，生成有序排查清单与历史对照"),
    ]
    for i, (title, trigger, output) in enumerate(scenarios):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.7), Inches(3.9), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(0.5),
                    title, font_size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(2.55), Inches(3.5), Inches(0.4),
                    f"触发：{trigger}", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(3.0),
                    output, font_size=14, color=LIGHT)


def slide_16_solution1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "方案一：知识库与学习系统", 16)
    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(12), Inches(0.45),
                "目标：把 20 年工艺沉淀转化为可检索、可传导、可生长的知识网络", font_size=16, color=LIGHT)
    add_bullet_textbox(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5), [
        "多源采集：工艺卡片、ELN、专利、文献、谱图",
        "结构化入库：分段切片 + 向量嵌入 + 图谱成网",
        "知识图谱可视化：节点浏览、边关系、谱图挂接",
        "对话式检索：RAG + meta-path 扩散 Top-K",
        "被动学习 + 主动学习双模式",
    ], font_size=14)
    add_textbox(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(0.4),
                "湃肽映射示例", font_size=16, color=ACCENT, bold=True)
    add_bullet_textbox(slide, Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.0), [
        "环肽纯化工艺卡片 → document 节点",
        "「聚集物控制」段落 → section + concept",
        "HPLC 色谱图 → media 节点",
        "批次实验 → experiment（序列+参数+结果）",
        "相似环肽工艺 → similar_to 长程联结",
    ], font_size=13, color=LIGHT)
    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.35),
                "验收：500 份工艺卡片入库，图谱可浏览，对话检索 Top-3 命中", font_size=13, color=ACCENT)


def slide_17_solution2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "方案二：智能选型决策引擎", 17)
    add_bullet_textbox(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(2.0), [
        "在知识网之上构建面向研发决策的智能体层，将「检索」升级为「推荐」",
        "输入目标分子特征 → PPR 沿 meta-path 扩散 → 汇聚相似实验与工艺 → 生成排序推荐",
    ], font_size=15)
    headers = ["智能体", "输入", "输出", "meta-path"]
    rows = [
        ["纯化选型", "氨基酸序列", "柱型/梯度/流速推荐", "molecule→similar_to→experiment"],
        ["合成路线", "肽长度/环化类型", "路线对比与保护基策略", "concept→section←experiment"],
        ["失效诊断", "收率/杂质异常", "排查清单+历史对照", "experiment→failed_with→concept"],
    ]
    add_table(slide, Inches(0.6), Inches(3.5), Inches(12.1), Inches(2.8),
              headers, rows, col_widths=[Inches(2.0), Inches(2.5), Inches(3.5), Inches(4.1)])
    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.35),
                "每条推荐附来源节点与专家审核入口，采纳率目标 > 60%", font_size=13, color=ACCENT)


def slide_18_model_strategy(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "模型选型：RAG 优先，Skill 调度", 18)
    strategies = [
        ("RAG 优先", "知识网约束检索 + 生成，不盲目微调大模型"),
        ("小模型专精", "OCR / 谱图描述 / 实体抽取用本地小模型"),
        ("Skill 调度", "按场景路由：检索 Skill、推荐 Skill、诊断 Skill"),
        ("不重复造轮", "底层模型选型开放，聚焦本体与传导引擎"),
    ]
    for i, (title, desc) in enumerate(strategies):
        y = Inches(1.55 + i * 1.35)
        add_textbox(slide, Inches(0.6), y, Inches(3.5), Inches(0.45),
                    title, font_size=18, color=ACCENT, bold=True)
        add_textbox(slide, Inches(4.3), y, Inches(8.3), Inches(0.9),
                    desc, font_size=15, color=LIGHT)
    add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.0),
                "原则：POC 阶段验证知识网价值，而非验证某个 LLM 的参数量\n"
                "幻觉防控：仅沿 meta-path 扩散 + 强制来源引用 + 专家审核门控",
                font_size=14, color=MUTED)


def slide_19_paradigm(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "范式跃迁：1.0 → 2.0 → 3.0", 19)
    eras = [
        ("1.0  经验驱动", "专家直觉 · 口口相传", "快但不可复制"),
        ("2.0  数据感知", "LIMS/ELN · 信息化", "可查但不成网"),
        ("3.0  智能研发", "知识网 + 智能体 + 闭环", "可传导、可行动、可生长"),
    ]
    for i, (era, feat, limit) in enumerate(eras):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.8), Inches(3.9), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT if i == 2 else MUTED
        add_textbox(slide, x + Inches(0.2), Inches(2.0), Inches(3.5), Inches(0.5),
                    era, font_size=18, color=ACCENT if i == 2 else WHITE, bold=True,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(0.8),
                    feat, font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.7), Inches(3.5), Inches(0.8),
                    limit, font_size=13, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            add_textbox(slide, x + Inches(3.5), Inches(3.0), Inches(0.6), Inches(0.4),
                        "→", font_size=28, color=ACCENT)
    add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
                "谛听方案推动湃肽从 2.0 迈向 3.0：让数据真正服务于研发决策", font_size=15, color=ACCENT,
                align=PP_ALIGN.CENTER)


def slide_20_wet_dry_loop(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "干湿闭环：实验驱动的知识进化", 20)
    steps = [
        ("01", "智能体推荐", "纯化/路线/诊断建议"),
        ("02", "专家审核", "采纳 / 修改 / 拒绝"),
        ("03", "实验执行", "湿实验验证方案"),
        ("04", "数据回流", "结果写入 experiment 节点"),
        ("05", "突触可塑", "边权更新，网络强化"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(0.65), Inches(2.0),
                                        Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        circle.line.color.rgb = ACCENT
        add_textbox(slide, x + Inches(0.65), Inches(2.3), Inches(1.2), Inches(0.6),
                    num, font_size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.4), Inches(2.4), Inches(0.4),
                    title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.85), Inches(2.4), Inches(0.8),
                    desc, font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(slide, x + Inches(2.1), Inches(2.4), Inches(0.4), Inches(0.4),
                        "→", font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5),
                "数字孪生进化：每一次湿实验结果都反哺干模型\n"
                "leads_to 边强化成功路径，failed_with 边标记失败警示\n"
                "知识网像神经系统一样「越用越灵敏」",
                font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)


def slide_21_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "实施路径：POC → 试点 → 推广", 21)
    headers = ["阶段", "周期", "范围", "验收标准"]
    rows = [
        ["POC", "8–10 周", "1 个品种 + 500 份工艺卡片",
         "图谱可浏览；对话检索 Top-3 命中"],
        ["试点", "3–6 月", "纯化选型智能体 + 3 条 meta-path",
         "新人培训可用；专家采纳率 > 60%"],
        ["推广", "6–12 月", "全管线 + 实验数据自动回流",
         "干湿闭环；边权可塑更新"],
    ]
    add_table(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(3.0),
              headers, rows, col_widths=[Inches(1.5), Inches(1.8), Inches(4.5), Inches(4.3)])
    add_bullet_textbox(slide, Inches(0.6), Inches(4.9), Inches(12), Inches(2.0), [
        "最小风险原则：POC 单品种验证，不一次性铺开",
        "数据主权：私有化部署，湃肽数据不出域",
        "每个阶段有明确验收指标，可按需暂停或加速",
    ], font_size=14, color=LIGHT)


def slide_22_cooperation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "合作模式", 22)
    modes = [
        ("联合共建", "双方投入领域专家与工程团队，共建湃肽专属知识网本体"),
        ("数据主权", "湃肽保留全部数据所有权；谛听提供技术与平台能力"),
        ("分期交付", "按 POC → 试点 → 推广分期签约，每阶段独立验收"),
        ("私有化部署", "cGMP/FDA 合规要求下，系统部署在湃肽内网环境"),
        ("持续运营", "联合运营小组：知识入库审核、meta-path 迭代、智能体优化"),
    ]
    for i, (title, desc) in enumerate(modes):
        y = Inches(1.5 + i * 1.1)
        add_textbox(slide, Inches(0.6), y, Inches(3.0), Inches(0.4),
                    title, font_size=16, color=ACCENT, bold=True)
        add_textbox(slide, Inches(3.8), y, Inches(8.8), Inches(0.8),
                    desc, font_size=14, color=LIGHT)


def slide_23_palantir(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "从 Palantir 借鉴的三件事", 23)
    lessons = [
        ("01  先建本体，再谈 AI",
         "Ontology 是业务语义层，不是数据库表。知识网本体是谛听方案的灵魂。"),
        ("02  让 AI 连接流程，而非替代流程",
         "AI 的价值在于 Action：推荐参数、生成方案、驱动实验，而非聊天。"),
        ("03  操作系统思维",
         "不做单点工具，构建可生长、可扩展、可传承的研发知识操作系统。"),
    ]
    for i, (title, desc) in enumerate(lessons):
        y = Inches(1.6 + i * 1.75)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), y, Inches(12.1), Inches(1.45))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x12, 0x22, 0x38)
        box.line.color.rgb = ACCENT
        add_textbox(slide, Inches(0.9), y + Inches(0.2), Inches(11.5), Inches(0.45),
                    title, font_size=18, color=ACCENT, bold=True)
        add_textbox(slide, Inches(0.9), y + Inches(0.7), Inches(11.5), Inches(0.6),
                    desc, font_size=14, color=LIGHT)


def slide_24_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(2.5), Inches(2.8), Inches(8.3))
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.9),
                "谛听  DITING", font_size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.9),
                "让每一条多肽工艺都被「看见」", font_size=30, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.6),
                "生物-数字双模态知识网 × 神经扩散智能体", font_size=18, color=LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
                "期待与湃肽生物共建多肽研发智能知识网络", font_size=16, color=MUTED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
                "谢谢 · THANK YOU", font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ── 构建入口 ──────────────────────────────────────────────
def build() -> Path:
    """生成完整 24 页 PPTX 并返回输出路径。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_partner(prs)
    slide_03_blackbox(prs)
    slide_04_silos(prs)
    slide_05_positioning(prs)
    slide_06_philosophy(prs)
    slide_07_comparison(prs)
    slide_08_capabilities(prs)
    slide_09_architecture(prs)
    slide_10_ontology_vs_model(prs)
    slide_11_flywheel(prs)
    slide_12_three_components(prs)
    slide_13_why_ontology(prs)
    slide_14_moat(prs)
    slide_15_scenarios(prs)
    slide_16_solution1(prs)
    slide_17_solution2(prs)
    slide_18_model_strategy(prs)
    slide_19_paradigm(prs)
    slide_20_wet_dry_loop(prs)
    slide_21_roadmap(prs)
    slide_22_cooperation(prs)
    slide_23_palantir(prs)
    slide_24_closing(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}  ({len(prs.slides)} slides)")
    return OUTPUT


if __name__ == "__main__":
    build()
